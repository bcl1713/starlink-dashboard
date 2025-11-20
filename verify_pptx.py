import io
import sys
from datetime import datetime, timedelta, timezone
from app.mission.models import Mission, MissionTimeline, TimelineSegment, TimelineStatus, Transport, TransportState, TransportConfig
from app.mission.exporter import generate_pptx_export, generate_pdf_export, TimelineExportFormat

# Mock data
def create_mock_timeline():
    # Create mock segments (11 segments to trigger orphan logic: 10+1 -> 8+3)
    segments = []
    start_time = datetime.now(timezone.utc)
    for i in range(11):
        segments.append(
            TimelineSegment(
                id=f"seg-{i}",
                segment_id=i + 1,
                start_time=start_time + timedelta(minutes=i*30),
                end_time=start_time + timedelta(minutes=(i+1)*30),
                duration=timedelta(minutes=30),
                status=TimelineStatus.NOMINAL if i % 2 == 0 else TimelineStatus.DEGRADED,
                x_transport_state=TransportState.AVAILABLE,
                ka_transport_state=TransportState.AVAILABLE,
                ku_transport_state=TransportState.AVAILABLE,
                reasons=["Safety-of-Flight"] if i == 5 else (["AAR"] if i == 6 else (["Test segment"] if i % 5 == 0 else []))
            )
        )
    
    timeline = MissionTimeline(
        mission_id="TEST-MISSION",
        created_at=datetime.now(),
        segments=segments,
        advisories=[],
        statistics={}
    )
    return timeline

def create_mock_mission():
    return Mission(
        id="TEST-MISSION",
        name="Test Mission",
        route_id="test-route",
        transports=TransportConfig(
            initial_x_satellite_id="X-1"
        )
    )

def verify_pptx():
    print("Verifying PPTX export...")
    timeline = create_mock_timeline()
    mission = create_mock_mission()
    
    try:
        # Generate PPTX
        pptx_bytes = generate_pptx_export(timeline, mission)
        print(f"Generated artifact: application/vnd.openxmlformats-officedocument.presentationml.presentation, size: {len(pptx_bytes)} bytes")
        
        if len(pptx_bytes) == 0:
            print("ERROR: Generated content is empty")
            return False
            
        # Try to load with python-pptx to verify validity
        from pptx import Presentation
        stream = io.BytesIO(pptx_bytes)
        prs = Presentation(stream)
        # Verify content
        print(f"Loaded presentation with {len(prs.slides)} slides")
        if len(prs.slides) != 3:
            print(f"ERROR: Expected 3 slides, got {len(prs.slides)}")
            sys.exit(1) # Changed from return False
            
        print("SUCCESS: PPTX export verified")
        
        # Verify PDF export
        print("Verifying PDF export...")
        try:
            pdf_bytes = generate_pdf_export(timeline, mission)
            print(f"Generated artifact: application/pdf, size: {len(pdf_bytes)} bytes")
            if len(pdf_bytes) > 0:
                 print("SUCCESS: PDF export verified")
            else:
                 print("FAILURE: PDF export empty")
                 sys.exit(1)
        except Exception as e:
            print(f"FAILURE: PDF export failed: {e}")
            import traceback
            traceback.print_exc()
            sys.exit(1)

        print("SUCCESS: All exports verified")
        sys.exit(0) # Changed from return True
        
    except Exception as e:
        print(f"ERROR: Verification failed: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1) # Changed from return False

if __name__ == "__main__":
    success = verify_pptx()
    sys.exit(0 if success else 1)
