"""Reusable PowerPoint presentation builder for mission exports.

This module consolidates PPTX generation logic previously duplicated across
`exporter/__main__.py` and `package/__main__.py`, providing a single source
of truth for slide creation, pagination, and styling.

Functions:
    create_pptx_presentation: Main entry point for generating a complete presentation
    add_route_map_slide: Add a slide with the mission route map
    add_timeline_table_slides: Add paginated timeline table slides
"""

from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import TYPE_CHECKING

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.mission.exporter.formatting import mission_start_timestamp
from app.mission.exporter.pptx_styling import (
    STATUS_CRITICAL,
    STATUS_DEGRADED,
    STATUS_NOMINAL,
    STATUS_SOF,
    TEXT_WHITE,
    add_footer_bar,
    add_footer_text,
    add_header_bar,
    add_logo,
    add_slide_title,
)
from app.mission.exporter.transport_utils import STATE_COLUMNS, TRANSPORT_DISPLAY
from app.mission.models import Transport

if TYPE_CHECKING:
    from pptx.slide import Slide

    from app.mission.models import Mission, MissionLeg
    from app.mission.exporter import MissionLegTimeline
    from app.routing import POIManager, RouteManager

logger = logging.getLogger(__name__)


def add_mission_slides_to_presentation(
    prs: Presentation,
    mission: Mission | MissionLeg | None,
    timeline: MissionLegTimeline,
    parent_mission_id: str | None = None,
    route_manager: RouteManager | None = None,
    poi_manager: POIManager | None = None,
    logo_path: Path | None = None,
    map_cache: dict[str, bytes] | None = None,
) -> None:
    """Add mission slides (route map and timeline tables) to an existing presentation.

    This function adds slides directly to the provided presentation object:
    - Route map slide with optional caching
    - Paginated timeline table slides with status coloring

    Args:
        prs: Presentation object to add slides to
        mission: Mission or MissionLeg object (None for standalone timeline)
        timeline: Timeline data for this mission/leg
        parent_mission_id: Parent mission ID (for multi-leg exports)
        route_manager: Route manager for map generation
        poi_manager: POI manager for map markers
        logo_path: Path to logo image file
        map_cache: Optional cache for generated maps (route_id -> bytes)
    """
    # Import _generate_route_map here to avoid circular imports
    from app.mission.exporter import _generate_route_map

    # Add route map slide
    add_route_map_slide(
        prs=prs,
        timeline=timeline,
        mission=mission,
        parent_mission_id=parent_mission_id,
        route_manager=route_manager,
        poi_manager=poi_manager,
        logo_path=logo_path,
        map_cache=map_cache,
        _generate_route_map=_generate_route_map,
    )

    # Add timeline table slides
    add_timeline_table_slides(
        prs=prs,
        timeline=timeline,
        mission=mission,
        logo_path=logo_path,
    )


def add_route_map_slide(
    prs: Presentation,
    timeline: MissionLegTimeline,
    mission: Mission | MissionLeg | None,
    parent_mission_id: str | None,
    route_manager: RouteManager | None,
    poi_manager: POIManager | None,
    logo_path: Path | None,
    map_cache: dict[str, bytes] | None,
    _generate_route_map,  # Injected to avoid circular import
) -> None:
    """Add route map slide to presentation.

    Args:
        prs: Presentation object to add slide to
        timeline: Timeline data for the route
        mission: Mission or MissionLeg object
        parent_mission_id: Parent mission ID (for multi-leg exports)
        route_manager: Route manager for map generation
        poi_manager: POI manager for map markers
        logo_path: Path to logo image
        map_cache: Optional cache for generated maps (route_id -> bytes)
        _generate_route_map: Map generation function (injected)
    """
    blank_slide_layout = prs.slide_layouts[6]
    slide_map = prs.slides.add_slide(blank_slide_layout)

    # Add header and footer bars
    add_header_bar(slide_map, 0, 0, 10, 0.15)
    add_footer_bar(slide_map, 0, 5.47, 10, 0.15)

    # Add logo
    if logo_path:
        add_logo(slide_map, logo_path, 0.2, 0.02, 0.5, 0.5)

    # Add slide title
    leg_name = timeline.mission_leg_id if timeline else "Route"
    add_slide_title(slide_map, f"{leg_name} - Route Map", top=0.2)

    # Get mission metadata for footer
    mission_id = mission.id if mission else timeline.mission_leg_id
    organization = (
        mission.description if (mission and mission.description) else "Organization"
    )

    # Generate map image (with caching support)
    route_id = mission.route_id if mission else None
    map_image_bytes = None

    # Check cache first
    if route_id and map_cache is not None and route_id in map_cache:
        map_image_bytes = map_cache[route_id]
        logger.info(f"Cache hit for route {route_id}")
    else:
        # Generate map
        try:
            map_image_bytes = _generate_route_map(
                timeline,
                mission,
                parent_mission_id=parent_mission_id,
                route_manager=route_manager,
                poi_manager=poi_manager,
            )
            logger.info(f"Cache miss for route {route_id}, generated map")

            # Store in cache if available
            if route_id and map_cache is not None:
                map_cache[route_id] = map_image_bytes
        except Exception as e:
            logger.error("Failed to generate map: %s", e, exc_info=True)

    # Add map to slide
    if map_image_bytes:
        try:
            map_image_stream = io.BytesIO(map_image_bytes)

            # Add image to slide, centered
            # Slide size is 10x5.62 inches
            # Available space: 10.0 wide x ~4.5 tall (below title, above footer)
            # Constrain height to 4.0 inches and center horizontally
            height = Inches(4.0)

            # Add picture first to get its dimensions after aspect ratio is applied
            pic = slide_map.shapes.add_picture(
                map_image_stream, Inches(0), Inches(0.9), height=height
            )

            # Center the picture horizontally
            pic.left = int((Inches(10) - pic.width) / 2)

        except Exception as e:
            logger.error("Failed to add map to PPTX: %s", e, exc_info=True)
            textbox = slide_map.shapes.add_textbox(
                Inches(1), Inches(1), Inches(8), Inches(1)
            )
            textbox.text = f"Map generation failed: {str(e)}"

    # Add footer text (centered within gold bar with white text)
    add_footer_text(
        slide_map,
        f"{mission_id} | {organization}",
        bottom=5.45,
        font_size=7,
        color=TEXT_WHITE,
    )


def add_timeline_table_slides(
    prs: Presentation,
    timeline: MissionLegTimeline,
    mission: Mission | MissionLeg | None,
    logo_path: Path | None,
) -> None:
    """Add paginated timeline table slides to presentation.

    This function creates one or more slides with timeline data, automatically
    paginating rows to prevent overcrowding and handling edge cases like orphan rows.

    Args:
        prs: Presentation object to add slides to
        timeline: Timeline data to display
        mission: Mission or MissionLeg object
        logo_path: Path to logo image
    """
    # Import here to avoid circular dependency
    from app.mission.exporter import _segment_rows

    timeline_df = _segment_rows(timeline, mission)

    if timeline_df.empty:
        return

    # Define columns to show
    columns_to_show = [
        "Segment #",
        "Status",
        "Start Time",
        "End Time",
        "Duration",
        TRANSPORT_DISPLAY[Transport.X],
        TRANSPORT_DISPLAY[Transport.KA],
        TRANSPORT_DISPLAY[Transport.KU],
        "Reasons",
    ]

    # Paginate timeline rows
    chunks = _paginate_timeline_rows(
        timeline_df, rows_per_slide=7, min_rows_last_slide=3
    )

    # Get mission metadata
    leg_name = timeline.mission_leg_id if timeline else "Mission"
    mission_id = mission.id if mission else timeline.mission_leg_id
    organization = (
        mission.description if (mission and mission.description) else "Organization"
    )

    # Get mission start date for footer
    mission_start = mission_start_timestamp(timeline)
    mission_date = mission_start.strftime("%Y-%m-%d")

    blank_slide_layout = prs.slide_layouts[6]

    for chunk_idx, chunk in enumerate(chunks):
        # Add new slide for each chunk
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add header and footer bars
        add_header_bar(slide, 0, 0, 10, 0.15)
        add_footer_bar(slide, 0, 5.47, 10, 0.15)

        # Add logo
        if logo_path:
            add_logo(slide, logo_path, 0.2, 0.02, 0.5, 0.5)

        # Add slide title
        title_text = (
            f"{leg_name} - Timeline"
            if chunk_idx == 0
            else f"{leg_name} - Timeline (cont.)"
        )
        add_slide_title(slide, title_text, top=0.2)

        # Add footer text
        footer_text = f"Date: {mission_date} | {mission_id} | {organization}"
        add_footer_text(slide, footer_text, bottom=5.45, font_size=7, color=TEXT_WHITE)

        # Add timeline table for this chunk
        _add_timeline_table(slide, chunk, columns_to_show)


def _paginate_timeline_rows(
    df: pd.DataFrame,
    rows_per_slide: int = 7,
    min_rows_last_slide: int = 3,
) -> list[pd.DataFrame]:
    """Split timeline rows into chunks for pagination.

    This function implements smart pagination that prevents orphan rows on the
    last page by moving rows from the second-to-last page when necessary.

    Args:
        df: DataFrame with timeline rows
        rows_per_slide: Maximum rows per slide (default: 7)
        min_rows_last_slide: Minimum rows for last slide (default: 3)

    Returns:
        List of DataFrame chunks, one per slide

    Examples:
        >>> df = pd.DataFrame({'col': range(11)})
        >>> chunks = _paginate_timeline_rows(df, rows_per_slide=10, min_rows_last_slide=3)
        >>> [len(c) for c in chunks]
        [8, 3]  # Prevents 10+1 split, creates 8+3 instead
    """
    total_rows = len(df)
    chunks = []

    if total_rows <= rows_per_slide:
        chunks.append(df)
    else:
        # Naive split first
        current_idx = 0
        while current_idx < total_rows:
            end_idx = min(current_idx + rows_per_slide, total_rows)
            chunks.append(df.iloc[current_idx:end_idx])
            current_idx = end_idx

        # Check last chunk and adjust if needed
        if len(chunks) > 1:
            last_chunk = chunks[-1]
            if len(last_chunk) < min_rows_last_slide:
                # Need to move items from second to last chunk
                # Re-build chunks: keep all but last two, then re-split the remainder
                base_chunks = chunks[:-2]
                remainder_count = len(chunks[-2]) + len(chunks[-1])

                # Calculate split point
                split_idx = (
                    total_rows - remainder_count
                )  # Start of second-to-last chunk
                second_last_len = remainder_count - min_rows_last_slide

                chunk_second_last = df.iloc[split_idx : split_idx + second_last_len]
                chunk_last = df.iloc[split_idx + second_last_len :]

                chunks = base_chunks + [chunk_second_last, chunk_last]

    return chunks


def _add_timeline_table(
    slide: Slide,
    chunk: pd.DataFrame,
    columns_to_show: list[str],
) -> None:
    """Add timeline table to slide with status coloring.

    This function creates a formatted table with:
    - Header row with dark blue background
    - Alternating row colors for readability
    - Status badge coloring (NOMINAL, SOF, DEGRADED, CRITICAL)
    - Transport state coloring (yellow/red for warnings)

    Args:
        slide: Slide object to add table to
        chunk: DataFrame chunk with timeline rows
        columns_to_show: List of column names to display
    """
    rows = len(chunk) + 1  # +1 for header
    cols = len(columns_to_show)

    # Create table
    left = Inches(0.5)
    top = Inches(0.9)
    width = Inches(9.0)
    # Use a minimal height so rows don't stretch to fill a large area
    height = Inches(1.0)

    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table

    # Set column widths (adjusted for wider times)
    # Total width 9 inches
    col_weights = [0.6, 1.0, 1.75, 1.75, 1.0, 1.0, 1.0, 1.0, 1.5]
    total_weight = sum(col_weights)
    for i, weight in enumerate(col_weights):
        table.columns[i].width = int(width * (weight / total_weight))

    # Header row style
    for i, col_name in enumerate(columns_to_show):
        cell = table.cell(0, i)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(51, 102, 178)  # Dark Blue
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White
        paragraph.font.bold = True
        paragraph.font.size = Pt(10)
        paragraph.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, (_, row_data) in enumerate(chunk.iterrows(), start=1):
        # Pre-calculate bad transports for this row
        bad_transports = []
        for col_name in STATE_COLUMNS:
            val = str(row_data[col_name] if row_data[col_name] else "").lower()
            if val in ("degraded", "warning", "offline"):
                bad_transports.append(col_name)

        is_critical_row = len(bad_transports) >= 2

        for col_idx, col_name in enumerate(columns_to_show):
            cell = table.cell(row_idx, col_idx)
            val = str(row_data[col_name] if row_data[col_name] is not None else "")

            # Override Status text if critical
            if col_name == "Status" and is_critical_row:
                val = "CRITICAL"

            cell.text = val

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.alignment = PP_ALIGN.LEFT

            # Alternating row colors
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light Gray
            else:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White

            # Status coloring for Status column with new color palette
            if col_name == "Status":
                val_lower = val.lower()

                # Check for Safety-of-Flight in reasons
                is_sof = False
                reasons = str(row_data.get("Reasons", "")).lower()
                if "safety-of-flight" in reasons or "aar" in reasons:
                    is_sof = True

                # Override Status text to "SOF" if it's a safety window
                if is_sof and val_lower in ("available", "nominal", "warning"):
                    cell.text = "SOF"
                    # Re-apply font size since setting text resets it
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(9)
                        paragraph.alignment = PP_ALIGN.LEFT

                # Apply status badge colors
                if is_critical_row and not is_sof:
                    # Critical status
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = STATUS_CRITICAL
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
                        paragraph.font.bold = True
                elif is_sof:
                    # Safety of Flight
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = STATUS_SOF
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
                        paragraph.font.bold = True
                elif val_lower in ("degraded", "warning"):
                    # Degraded status
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = STATUS_DEGRADED
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
                        paragraph.font.bold = True
                elif val_lower in ("nominal", "available"):
                    # Nominal status
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = STATUS_NOMINAL
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
                        paragraph.font.bold = True

            # Keep existing coloring for transport state columns
            elif col_name in STATE_COLUMNS:
                val_lower = val.lower()
                if val_lower in ("degraded", "warning", "offline"):
                    cell.fill.solid()
                    if val_lower in ("degraded", "warning"):
                        cell.fill.fore_color.rgb = RGBColor(
                            255, 255, 204
                        )  # Light Yellow
                    else:
                        cell.fill.fore_color.rgb = RGBColor(255, 204, 204)  # Light Red
