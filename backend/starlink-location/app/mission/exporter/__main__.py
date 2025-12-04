"""Mission timeline export utilities.

Transforms `MissionLegTimeline` data into CSV, XLSX, and PDF deliverables with
parallel timestamp formats (UTC, Eastern, and T+ offsets) suitable for
customer-facing mission briefs.
"""

# FR-004: File exceeds 300 lines (2126 lines) because mission export generation
# requires complex multi-format handling (PDF graphics with matplotlib, XLSX with
# formulas, CSV with computed fields) that cannot be cleanly separated. Refactoring
# would require extraction into 4+ service modules with circular dependencies on
# timeline builders. Deferred to v0.4.0 pending SOLID improvements.

from __future__ import annotations

import io
import json
import logging
import matplotlib

matplotlib.use("Agg")  # Headless mode for Docker
import matplotlib.pyplot as plt
import cartopy.crs as ccrs
import cartopy.feature as cfeature
from matplotlib.lines import Line2D
from dataclasses import dataclass
from datetime import datetime, timezone
from enum import Enum
from pathlib import Path
from zoneinfo import ZoneInfo

import pandas as pd
import openpyxl
from openpyxl.drawing.image import Image as OpenpyxlImage
from openpyxl.styles import PatternFill
from reportlab.lib import colors
from reportlab.lib.enums import TA_RIGHT
from reportlab.lib.pagesizes import LETTER, landscape
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    Image,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
    PageBreak,
)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.mission.models import (
    Mission,
    MissionLegTimeline,
    TimelineSegment,
    TimelineStatus,
    Transport,
)
from app.mission.exporter.formatting import (
    compose_time_block,
    ensure_timezone,
    format_eastern,
    format_offset,
    format_seconds_hms,
    format_utc,
    humanize_metric_name,
    mission_start_timestamp,
)
from app.mission.exporter.transport_utils import (
    TRANSPORT_DISPLAY,
    STATE_COLUMNS,
    STATUS_COLORS,
    display_transport_state,
    segment_is_x_ku_warning,
    serialize_transport_list,
)
from app.services.poi_manager import POIManager
from app.services.route_manager import RouteManager

logger = logging.getLogger(__name__)

EASTERN_TZ = ZoneInfo("America/New_York")
LIGHT_YELLOW = colors.Color(1.0, 1.0, 0.85)
LIGHT_RED = colors.Color(1.0, 0.85, 0.85)
LOGO_PATH = Path(__file__).parent.with_name("assets").joinpath("logo.png")


def _load_logo_flowable() -> Image | None:
    """Return a scaled logo image if available."""
    if not LOGO_PATH.exists():
        return None
    try:
        logo = Image(str(LOGO_PATH))
    except OSError:
        return None
    max_width = 1.6 * inch
    max_height = 0.9 * inch
    logo._restrictSize(max_width, max_height)
    logo.hAlign = "RIGHT"
    return logo


class ExportGenerationError(RuntimeError):
    """Raised when timeline export generation fails."""


class TimelineExportFormat(str, Enum):
    """Supported export formats."""

    CSV = "csv"
    XLSX = "xlsx"
    PDF = "pdf"
    PPTX = "pptx"

    @classmethod
    def from_string(cls, raw: str) -> "TimelineExportFormat":
        """Parse user-supplied format strings case-insensitively."""
        value = (raw or "").strip().lower()
        try:
            return cls(value)
        except ValueError as exc:
            raise ExportGenerationError(f"Unsupported export format: {raw}") from exc


@dataclass(slots=True)
class ExportArtifact:
    """Container describing a generated export file."""

    content: bytes
    media_type: str
    extension: str


# Utility functions imported from exporter modules above
# - ensure_timezone, format_utc, format_eastern, format_offset from formatting
# - serialize_transport_list, is_x_ku_conflict_reason, segment_is_x_ku_warning,
#   display_transport_state from transport_utils


def _aar_block_rows(
    timeline: MissionLegTimeline,
    mission: Mission | None,
    mission_start: datetime,
) -> list[tuple[datetime, int, dict]]:
    blocks = timeline.statistics.get("_aar_blocks") if timeline.statistics else None
    if not blocks:
        return []
    rows: list[tuple[datetime, int, dict]] = []
    for block in blocks:
        start_raw = block.get("start")
        end_raw = block.get("end")
        if not (start_raw and end_raw):
            continue
        start_dt = _parse_iso_timestamp(start_raw)
        end_dt = _parse_iso_timestamp(end_raw)
        if end_dt <= start_dt:
            continue
        rows.append(
            (
                start_dt,
                0,
                _build_aar_record(
                    start_dt,
                    end_dt,
                    mission,
                    timeline,
                    mission_start,
                ),
            )
        )
    return rows


def _parse_iso_timestamp(raw: str) -> datetime:
    try:
        dt = datetime.fromisoformat(raw)
    except ValueError:
        dt = datetime.fromtimestamp(0, tz=timezone.utc)
    return ensure_timezone(dt)


def _build_aar_record(
    start: datetime,
    end: datetime,
    mission: Mission | None,
    timeline: MissionLegTimeline,
    mission_start: datetime,
) -> dict:
    duration = max((end - start).total_seconds(), 0)
    segment = _segment_at_time(timeline, start)
    x_state = display_transport_state(segment.x_state) if segment else ""
    ka_state = display_transport_state(segment.ka_state) if segment else ""
    ku_state = display_transport_state(segment.ku_state) if segment else ""
    return {
        "Segment #": "AAR",
        "Mission ID": mission.id if mission else timeline.mission_id,
        "Mission Name": (
            mission.name if mission and mission.name else timeline.mission_id
        ),
        "Status": "WARNING",
        "Start Time": compose_time_block(start, mission_start),
        "End Time": compose_time_block(end, mission_start),
        "Duration": format_seconds_hms(duration),
        TRANSPORT_DISPLAY[Transport.X]: x_state,
        TRANSPORT_DISPLAY[Transport.KA]: ka_state,
        TRANSPORT_DISPLAY[Transport.KU]: ku_state,
        "Impacted Transports": "",
        "Reasons": "AAR",
        "Metadata": "",
    }


def _segment_at_time(
    timeline: MissionLegTimeline,
    timestamp: datetime,
) -> TimelineSegment | None:
    ts = ensure_timezone(timestamp)
    for segment in timeline.segments:
        start = ensure_timezone(segment.start_time)
        end = (
            ensure_timezone(segment.end_time)
            if segment.end_time
            else datetime.max.replace(tzinfo=timezone.utc)
        )
        if start <= ts < end:
            return segment
    return timeline.segments[-1] if timeline.segments else None


def _get_detailed_segment_statuses(
    start_time: datetime, end_time: datetime, timeline: MissionLegTimeline
) -> list[tuple[datetime, datetime, str]]:
    if not timeline or not timeline.segments:
        return [(start_time, end_time, "unknown")]

    intervals = []
    current_time = start_time

    # Optimization: Filter relevant segments first
    relevant_segments = []
    for seg in timeline.segments:
        s_start = ensure_timezone(seg.start_time)
        s_end = (
            ensure_timezone(seg.end_time)
            if seg.end_time
            else datetime.max.replace(tzinfo=timezone.utc)
        )
        if s_end > start_time and s_start < end_time:
            relevant_segments.append((s_start, s_end, seg))

    # Sort by start time
    relevant_segments.sort(key=lambda x: x[0])

    # Also check for AAR blocks
    aar_blocks = []
    if timeline.statistics and "_aar_blocks" in timeline.statistics:
        for block in timeline.statistics["_aar_blocks"]:
            aar_start_raw = block.get("start")
            aar_end_raw = block.get("end")
            if aar_start_raw and aar_end_raw:
                aar_start = _parse_iso_timestamp(aar_start_raw)
                aar_end = _parse_iso_timestamp(aar_end_raw)
                if aar_end > start_time and aar_start < end_time:
                    aar_blocks.append((aar_start, aar_end))

    while current_time < end_time:
        # Find the segment that covers current_time
        active_seg = None
        next_change = end_time

        for s_start, s_end, seg in relevant_segments:
            if s_start <= current_time < s_end:
                active_seg = seg
                next_change = min(end_time, s_end)
                break
            elif s_start > current_time:
                # This is a future segment, it might define the end of a gap
                next_change = min(end_time, s_start)
                break

        # Check if current_time falls within an AAR block
        in_aar_block = False
        for aar_start, aar_end in aar_blocks:
            if aar_start <= current_time < aar_end:
                in_aar_block = True
                # AAR block might end before the segment ends
                next_change = min(next_change, aar_end)
                break
            elif aar_start > current_time:
                # Future AAR block might start before next_change
                next_change = min(next_change, aar_start)

        # Determine status
        status = "nominal"  # Default if no segment found (gap) or segment has no status
        if active_seg:
            raw_status = active_seg.status.value
            reasons = str(active_seg.reasons).lower()
            is_sof = "safety-of-flight" in reasons or "aar" in reasons

            if raw_status == "critical":
                status = "critical"
            elif raw_status == "degraded" or is_sof or in_aar_block:
                status = "degraded"
            elif raw_status == "warning":
                status = "degraded"
            else:
                status = "nominal"
        else:
            status = "unknown"  # Gap in timeline

        # Override to degraded if in AAR block (even if segment is nominal)
        if in_aar_block and status == "nominal":
            status = "degraded"

        intervals.append((current_time, next_change, status))
        current_time = next_change

    return intervals


def _interpolate_position_at_time(target_time, p1, p2):
    if not p1.expected_arrival_time or not p2.expected_arrival_time:
        return p1

    t1 = ensure_timezone(p1.expected_arrival_time)
    t2 = ensure_timezone(p2.expected_arrival_time)

    if target_time <= t1:
        return p1
    if target_time >= t2:
        return p2

    total_duration = (t2 - t1).total_seconds()
    if total_duration == 0:
        return p1

    elapsed = (target_time - t1).total_seconds()
    fraction = elapsed / total_duration

    lat = p1.latitude + (p2.latitude - p1.latitude) * fraction

    # Handle IDL for Longitude
    lon1 = p1.longitude
    lon2 = p2.longitude

    if abs(lon2 - lon1) > 180:
        # Crosses IDL
        if lon1 < 0:
            lon1 += 360
        if lon2 < 0:
            lon2 += 360

        lon = lon1 + (lon2 - lon1) * fraction
        if lon > 180:
            lon -= 360
    else:
        lon = lon1 + (lon2 - lon1) * fraction

    return type("Point", (), {"latitude": lat, "longitude": lon})


def _generate_route_map(
    timeline: MissionLegTimeline,
    mission: Mission | None = None,
    parent_mission_id: str | None = None,
    route_manager: RouteManager | None = None,
    poi_manager: POIManager | None = None,
) -> bytes:
    """Generate a 4K PNG image of the route map.

    Current phase: Route line drawing (Phase 9)
    - Output: 3840x2160 pixels @ 300 DPI (16:9 4K)
    - Map bounds calculated from route waypoints
    - Route line drawn with IDL crossing handling (split segments)
    - Projection centers on route (handles Pacific view for IDL crossings)

    Args:
        timeline: The mission timeline with segments and timing data
        mission: The mission object containing route and POI information (optional)
        parent_mission_id: Parent mission ID when exporting legs from a multi-leg mission (optional)
        route_manager: RouteManager instance for fetching route data
        poi_manager: POIManager instance for fetching POI data

    Returns:
        PNG image as bytes.
    """
    # Phase 9: Draw route as simple line with IDL handling and Pacific centering

    # Canvas dimensions for 4K 16:9
    # 3840 x 2160 pixels
    pixel_width = 3840
    pixel_height = 2160
    dpi = 300
    width_inches = pixel_width / dpi  # 12.8 inches
    height_inches = pixel_height / dpi  # 7.2 inches

    # Check if we have valid mission and route data
    if mission is None:
        logger.warning("_generate_route_map: mission is None, returning base canvas")
        return _base_map_canvas()
    if not mission.route_id:
        logger.warning(
            f"_generate_route_map: mission {mission.id} has no route_id, returning base canvas"
        )
        return _base_map_canvas()
    if not route_manager:
        logger.warning(
            "_generate_route_map: route_manager is None, returning base canvas"
        )
        return _base_map_canvas()

    logger.info(
        f"_generate_route_map: Generating map for mission={mission.id}, route={mission.route_id}"
    )

    # Fetch route from manager
    route = route_manager.get_route(mission.route_id)
    if route is None:
        # Log available routes to help diagnose mismatch
        available_routes = route_manager.list_routes()
        logger.error(f"Route '{mission.route_id}' not found in RouteManager!")
        logger.error(f"Available routes: {list(available_routes.keys())}")
        logger.error(f"Mission: {mission.id}, Leg: {mission.name}")
        return _base_map_canvas()

    if not route.points:
        logger.warning(
            f"Route '{mission.route_id}' has no points, returning base canvas"
        )
        return _base_map_canvas()

    # Extract waypoint coordinates and filter invalid points
    raw_points = route.points
    valid_points = [
        p for p in raw_points if p.latitude is not None and p.longitude is not None
    ]

    if len(valid_points) < len(raw_points):
        logger.warning(
            f"Filtered {len(raw_points) - len(valid_points)} invalid points from route {mission.route_id}"
        )

    if not valid_points:
        logger.warning(
            f"No valid points found for route {mission.route_id}, returning base canvas"
        )
        return _base_map_canvas()

    lats = [p.latitude for p in valid_points]
    lons = [p.longitude for p in valid_points]

    # Assign valid_points to points for downstream usage
    points = valid_points

    # Debug logging
    logger.info(f"Map generation - Route has {len(points)} valid points")

    # Detect IDL (International Date Line) crossings
    idl_crossing_segments = set()
    for i in range(len(valid_points) - 1):
        p1, p2 = valid_points[i], valid_points[i + 1]
        lon_diff = abs(p2.longitude - p1.longitude)
        if lon_diff > 180:
            idl_crossing_segments.add(i)

    is_idl_crossing = bool(idl_crossing_segments)

    # Determine projection center
    if is_idl_crossing:
        # For IDL crossing routes, center on the route's midpoint in normalized space
        # Normalize longitudes to 0-360 range for bounds calculation relative to 180
        norm_lons = []
        for lon in lons:
            if lon < 0:
                norm_lons.append(lon + 360)
            else:
                norm_lons.append(lon)

        min_lon = min(norm_lons)
        max_lon = max(norm_lons)
        min_lat, max_lat = min(lats), max(lats)

        # Calculate center of the route
        central_longitude = (min_lon + max_lon) / 2

        # Ensure central_longitude is within -180 to 180 for PlateCarree
        if central_longitude > 180:
            central_longitude -= 360

        logger.info(f"IDL Crossing Route - Dynamic Center: {central_longitude:.2f}")
        logger.info(f"Normalized Lon Range: {min_lon:.2f} to {max_lon:.2f}")
    else:
        # Standard route
        min_lon, max_lon = min(lons), max(lons)
        min_lat, max_lat = min(lats), max(lats)

        # Center on the route midpoint
        central_longitude = (min_lon + max_lon) / 2
        logger.info(f"Standard Route - Central Lon: {central_longitude:.2f}")

    # Calculate route extent
    lat_range = max_lat - min_lat
    lon_range = max_lon - min_lon

    # Enforce minimum extent to prevent extreme zooming (which can cause memory errors)
    MIN_EXTENT = 10.0
    if lat_range < MIN_EXTENT:
        center_lat = (min_lat + max_lat) / 2
        min_lat = center_lat - (MIN_EXTENT / 2)
        max_lat = center_lat + (MIN_EXTENT / 2)
        lat_range = MIN_EXTENT

    if lon_range < MIN_EXTENT:
        center_lon = (min_lon + max_lon) / 2
        min_lon = center_lon - (MIN_EXTENT / 2)
        max_lon = center_lon + (MIN_EXTENT / 2)
        lon_range = MIN_EXTENT

    # Canvas aspect ratio: 3840 / 2160 = 1.778 (16:9)
    canvas_aspect = pixel_width / pixel_height

    # Calculate natural geographic bounds with basic 5% padding
    # Start by padding the larger dimension
    if lon_range >= lat_range:
        # Longitude is dominant - pad it
        padding_lon = lon_range * 0.05
        bounds_west = min_lon - padding_lon
        bounds_east = max_lon + padding_lon
        padded_lon_range = bounds_east - bounds_west

        # Latitude stays at route extent + small margin
        padding_lat = lat_range * 0.02  # Just 2% for latitude
        bounds_south = min_lat - padding_lat
        bounds_north = max_lat + padding_lat
        padded_lat_range = bounds_north - bounds_south
    else:
        # Latitude is dominant - pad it
        padding_lat = lat_range * 0.05
        bounds_south = min_lat - padding_lat
        bounds_north = max_lat + padding_lat
        padded_lat_range = bounds_north - bounds_south

        # Longitude stays at route extent + small margin
        padding_lon = lon_range * 0.02  # Just 2% for longitude
        bounds_west = min_lon - padding_lon
        bounds_east = max_lon + padding_lon
        padded_lon_range = bounds_east - bounds_west

    # Geographic aspect ratio (lon / lat)
    geographic_aspect = padded_lon_range / padded_lat_range

    logger.info(
        f"Canvas aspect: {canvas_aspect:.3f}, Geographic aspect: {geographic_aspect:.3f}"
    )

    # Adjust bounds to match canvas aspect ratio without stretching
    if geographic_aspect < canvas_aspect:
        # Geographic bounds are taller than canvas - add padding to longitude
        needed_lon_range = padded_lat_range * canvas_aspect
        lon_excess = needed_lon_range - padded_lon_range
        bounds_west -= lon_excess / 2
        bounds_east += lon_excess / 2
        logger.info(f"Added longitude padding: {lon_excess:.2f}°")
    elif geographic_aspect > canvas_aspect:
        # Geographic bounds are wider than canvas - add padding to latitude
        needed_lat_range = padded_lon_range / canvas_aspect
        lat_excess = needed_lat_range - padded_lat_range
        bounds_south -= lat_excess / 2
        bounds_north += lat_excess / 2
        logger.info(f"Added latitude padding: {lat_excess:.2f}°")

    # Clamp latitude to valid ranges
    bounds_south = max(-90, bounds_south)
    bounds_north = min(90, bounds_north)

    # For longitude, if we are in 0-360 space (IDL crossing), we might go beyond 0 or 360 with padding
    # We'll handle this by transforming back to -180/180 for set_extent if needed,
    # OR just pass the values if using a shifted projection.

    logger.info(
        f"Calculated Bounds: W={bounds_west:.2f}, E={bounds_east:.2f}, S={bounds_south:.2f}, N={bounds_north:.2f}"
    )

    # Create 4K figure with custom projection
    fig = plt.figure(figsize=(width_inches, height_inches), dpi=dpi, facecolor="white")

    # Use PlateCarree with the calculated central longitude
    projection = ccrs.PlateCarree(central_longitude=central_longitude)
    ax = fig.add_subplot(111, projection=projection)

    # Remove all padding and margins
    fig.subplots_adjust(left=0, right=1, top=1, bottom=0)

    # Set map extent
    # Note: set_extent expects coordinates in the CRS provided by the `crs` argument.
    # If we use the SAME projection as the axis, we can pass the bounds we calculated in that space.

    if is_idl_crossing:
        # For IDL crossing, our bounds are in normalized 0-360 space.
        # The projection is centered at central_longitude.
        # We need to shift the bounds to be relative to the central longitude for the extent.
        # e.g. if Center=180, and Bounds=[170, 190], Extent should be [-10, 10].
        extent_west = bounds_west - central_longitude
        extent_east = bounds_east - central_longitude

        ax.set_extent(
            [extent_west, extent_east, bounds_south, bounds_north], crs=projection
        )
    else:
        # Standard case - bounds are standard -180..180
        # Projection is centered on route midpoint
        # We also need to shift these bounds to be relative to the center?
        # Yes, if we pass crs=projection (which is centered at C), we should pass relative coords.
        # e.g. Center=0, Bounds=[-10, 10]. Extent=[-10, 10].
        # e.g. Center=10, Bounds=[0, 20]. Extent=[-10, 10].
        extent_west = bounds_west - central_longitude
        extent_east = bounds_east - central_longitude

        ax.set_extent(
            [extent_west, extent_east, bounds_south, bounds_north], crs=projection
        )

    # Add map features
    ax.coastlines(resolution="50m", linewidth=0.5, color="#2c3e50")
    ax.add_feature(cfeature.BORDERS, linewidth=0.5, color="#bdc3c7")
    ax.add_feature(cfeature.LAND, facecolor="#ecf0f1", edgecolor="none")
    ax.add_feature(cfeature.OCEAN, facecolor="#d5e8f7", edgecolor="none")

    # Subtle gridlines without labels
    ax.gridlines(
        draw_labels=False, alpha=0.1, linestyle="-", linewidth=0.3, color="#95a5a6"
    )

    # Remove axis ticks and spines for clean appearance
    ax.spines["geo"].set_visible(False)
    ax.set_xticks([])
    ax.set_yticks([])

    # Phase 10: Draw route with color-coded segments based on timeline status
    if points:
        # Default to unknown/gray if no timeline data
        default_color = STATUS_COLORS["unknown"]

        # For each route segment (between consecutive waypoints), determine its color
        for i in range(len(points) - 1):
            p1 = points[i]
            p2 = points[i + 1]

            # If we lack timing, fall back to proportional coloring based on timeline status distribution
            if not p1.expected_arrival_time or not p2.expected_arrival_time:
                # Use the overall timeline status distribution to color the route segments
                # Distribute timeline status across route proportionally
                if not timeline or not timeline.segments:
                    fallback_color = STATUS_COLORS["unknown"]
                else:
                    # Calculate proportions of timeline statuses
                    stats = (
                        timeline.statistics
                        if timeline.statistics and isinstance(timeline.statistics, dict)
                        else {}
                    )
                    critical_secs = stats.get("critical_seconds", 0)
                    degraded_secs = stats.get("degraded_seconds", 0)
                    nominal_secs = stats.get("nominal_seconds", 0)
                    total_secs = (
                        critical_secs + degraded_secs + nominal_secs or 1
                    )  # Avoid division by zero

                    # Distribute route segments by proportion
                    # Use cumulative distribution to assign colors proportionally
                    critical_portion = critical_secs / total_secs
                    degraded_portion = (critical_secs + degraded_secs) / total_secs
                    segment_index = i / (len(points) - 1)  # 0-1 normalized position

                    if segment_index < critical_portion:
                        fallback_color = STATUS_COLORS["critical"]
                    elif segment_index < degraded_portion:
                        fallback_color = STATUS_COLORS["degraded"]
                    else:
                        fallback_color = STATUS_COLORS["nominal"]

                ax.plot(
                    [p1.longitude, p2.longitude],
                    [p1.latitude, p2.latitude],
                    color=fallback_color,
                    linewidth=1.5,
                    transform=ccrs.PlateCarree(),
                    zorder=5,
                )
                continue

            t1 = ensure_timezone(p1.expected_arrival_time)
            t2 = ensure_timezone(p2.expected_arrival_time)

            # Get sub-segments based on status
            sub_segments = _get_detailed_segment_statuses(t1, t2, timeline)

            for sub_start, sub_end, status in sub_segments:
                # Interpolate positions
                sp1 = _interpolate_position_at_time(sub_start, p1, p2)
                sp2 = _interpolate_position_at_time(sub_end, p1, p2)

                color = STATUS_COLORS.get(status, default_color)

                # Draw this sub-segment
                # Check for IDL crossing within this sub-segment
                if abs(sp2.longitude - sp1.longitude) > 180:
                    # Handle International Date Line crossings: split route into segments at ±180° boundaries
                    d_lon_short_path = 360 - abs(sp1.longitude - sp2.longitude)

                    if d_lon_short_path == 0:
                        continue

                    fraction = (180 - abs(sp1.longitude)) / d_lon_short_path
                    lat_at_180 = sp1.latitude + (sp2.latitude - sp1.latitude) * fraction

                    # Segment 1: P1 to Meridian
                    target_lon1 = 180 if sp1.longitude > 0 else -180
                    ax.plot(
                        [sp1.longitude, target_lon1],
                        [sp1.latitude, lat_at_180],
                        color=color,
                        linewidth=1.5,
                        transform=ccrs.PlateCarree(),
                        zorder=5,
                    )

                    # Segment 2: Meridian to P2
                    target_lon2 = 180 if sp2.longitude > 0 else -180
                    ax.plot(
                        [target_lon2, sp2.longitude],
                        [lat_at_180, sp2.latitude],
                        color=color,
                        linewidth=1.5,
                        transform=ccrs.PlateCarree(),
                        zorder=5,
                    )
                else:
                    # Normal segment
                    ax.plot(
                        [sp1.longitude, sp2.longitude],
                        [sp1.latitude, sp2.latitude],
                        color=color,
                        linewidth=1.5,
                        transform=ccrs.PlateCarree(),
                        zorder=5,
                    )

            # Add rounded caps for smooth joins (optional, but matplotlib lines usually join well)
            # For individual segments, we might see gaps if linewidth is large.
            # Adding a point marker at p1 can help smooth it, but might be overkill.
            # With linewidth 6, simple plot is usually fine.

    # Phase 11: Add POI Markers (Start, End, Waypoints)
    if points:
        # 1. Resolve Labels for Departure and Arrival
        start_label = "DEP"
        end_label = "ARR"

        # Build a map of waypoint names to coordinates for easy lookup
        waypoint_map = {}
        if route.waypoints:
            for wp in route.waypoints:
                if wp.name:
                    waypoint_map[wp.name] = wp

                # Also resolve Dep/Arr labels
                if wp.role == "departure" and wp.name:
                    start_label = wp.name
                elif wp.role == "arrival" and wp.name:
                    end_label = wp.name

        # Helper to interpolate position from timestamp
        def interpolate_position(target_time):
            if not points:
                return None

            # Find surrounding points
            for i in range(len(points) - 1):
                p1 = points[i]
                p2 = points[i + 1]
                if not p1.expected_arrival_time or not p2.expected_arrival_time:
                    continue

                if p1.expected_arrival_time <= target_time <= p2.expected_arrival_time:
                    # Linear interpolation
                    total_duration = (
                        p2.expected_arrival_time - p1.expected_arrival_time
                    ).total_seconds()
                    if total_duration == 0:
                        return p1

                    elapsed = (target_time - p1.expected_arrival_time).total_seconds()
                    fraction = elapsed / total_duration

                    lat = p1.latitude + (p2.latitude - p1.latitude) * fraction
                    lon = p1.longitude + (p2.longitude - p1.longitude) * fraction

                    # Handle IDL crossing for interpolation if needed (simplified here)
                    if abs(p2.longitude - p1.longitude) > 180:
                        # If crossing IDL, simple linear interp on lon is wrong, but for POI placement
                        # on a fine-grained route, it's usually close enough or we can skip.
                        # For now, return p1 to be safe.
                        return p1

                    return type("Point", (), {"latitude": lat, "longitude": lon})
            return None

        # 2. Collect Mission Event POIs (AAR + Sat Swaps)
        mission_event_pois = []

        if poi_manager and mission:
            # Get POIs for this specific leg (route_id + mission_id)
            # This ensures we only show POIs for the current leg, not all legs in the mission
            # Use parent_mission_id if provided (for multi-leg mission exports), otherwise use mission.id
            effective_mission_id = (
                parent_mission_id if parent_mission_id else mission.id
            )

            if mission.route_id:
                pois = poi_manager.list_pois(
                    route_id=mission.route_id, mission_id=effective_mission_id
                )
            else:
                pois = poi_manager.list_pois(mission_id=effective_mission_id)

            logger.info(
                f"Map generation: Found {len(pois)} POIs for route={mission.route_id}, mission={effective_mission_id} (parent={parent_mission_id})"
            )
            for poi in pois:
                mission_event_pois.append(
                    {
                        "lat": poi.latitude,
                        "lon": poi.longitude,
                        "label": poi.name,
                        "type": "mission_event",
                    }
                )
                logger.debug(
                    f"  - Adding POI to map: {poi.name} at ({poi.latitude}, {poi.longitude})"
                )
        else:
            # Fallback to manual extraction if POI manager not available
            # AAR Waypoints - Explicitly label Start and End
            if mission and mission.transports and mission.transports.aar_windows:
                for aar in mission.transports.aar_windows:
                    # AAR Start
                    if aar.start_waypoint_name in waypoint_map:
                        wp = waypoint_map[aar.start_waypoint_name]
                        mission_event_pois.append(
                            {
                                "lat": wp.latitude,
                                "lon": wp.longitude,
                                "label": "AAR Start",
                                "type": "aar_start",
                            }
                        )

                    # AAR End
                    if aar.end_waypoint_name in waypoint_map:
                        wp = waypoint_map[aar.end_waypoint_name]
                        mission_event_pois.append(
                            {
                                "lat": wp.latitude,
                                "lon": wp.longitude,
                                "label": "AAR End",
                                "type": "aar_end",
                            }
                        )

            # Satellite Transitions (X Swaps) - Configured
            if mission and mission.transports and mission.transports.x_transitions:
                for trans in mission.transports.x_transitions:
                    # Create a label for the transition (e.g., "Swap to X-2")
                    label = f"Swap {trans.target_satellite_id}"
                    mission_event_pois.append(
                        {
                            "lat": trans.latitude,
                            "lon": trans.longitude,
                            "label": label,
                            "type": "transition",
                        }
                    )

            # Ka Transitions (Auto-calculated from Timeline)
            # Iterate segments to detect changes in Ka satellites
            if timeline and timeline.segments:
                current_ka = set()
                # Initialize with first segment's Ka state if available
                first_seg = timeline.segments[0]
                if (
                    first_seg.metadata
                    and "satellites" in first_seg.metadata
                    and "Ka" in first_seg.metadata["satellites"]
                ):
                    current_ka = set(first_seg.metadata["satellites"]["Ka"])

                for i in range(1, len(timeline.segments)):
                    seg = timeline.segments[i]
                    next_ka = set()
                    if (
                        seg.metadata
                        and "satellites" in seg.metadata
                        and "Ka" in seg.metadata["satellites"]
                    ):
                        next_ka = set(seg.metadata["satellites"]["Ka"])

                    # Check for change
                    if current_ka != next_ka:
                        # Transition detected at seg.start_time
                        pos = interpolate_position(seg.start_time)
                        if pos:
                            # Construct label: "Swap POR to AOR"
                            added = next_ka - current_ka
                            removed = current_ka - next_ka

                            label_parts = []
                            if removed:
                                label_parts.append(f"{'/'.join(removed)}")
                            if added:
                                label_parts.append(f"{'/'.join(added)}")

                            if label_parts:
                                label = f"Swap {' to '.join(label_parts)}"
                                mission_event_pois.append(
                                    {
                                        "lat": pos.latitude,
                                        "lon": pos.longitude,
                                        "label": label,
                                        "type": "ka_transition",
                                    }
                                )

                        current_ka = next_ka

        # 3. Plot Departure Point (Start)
        start_point = points[0]
        ax.plot(
            start_point.longitude,
            start_point.latitude,
            marker="o",
            color="#2ecc71",
            markersize=12,
            markeredgecolor="white",
            markeredgewidth=2,
            transform=ccrs.PlateCarree(),
            zorder=10,
        )

        # Label for Departure
        ax.text(
            start_point.longitude + 0.5,
            start_point.latitude + 0.5,
            start_label,
            transform=ccrs.PlateCarree(),
            fontsize=10,
            fontweight="bold",
            color="#2c3e50",
            zorder=11,
            bbox=dict(facecolor="white", alpha=0.7, edgecolor="none", pad=1),
        )

        # 4. Plot Arrival Point (End)
        end_point = points[-1]
        ax.plot(
            end_point.longitude,
            end_point.latitude,
            marker="o",
            color="#e74c3c",
            markersize=12,
            markeredgecolor="white",
            markeredgewidth=2,
            transform=ccrs.PlateCarree(),
            zorder=10,
        )

        # Label for Arrival
        ax.text(
            end_point.longitude + 0.5,
            end_point.latitude + 0.5,
            end_label,
            transform=ccrs.PlateCarree(),
            fontsize=10,
            fontweight="bold",
            color="#2c3e50",
            zorder=11,
            bbox=dict(facecolor="white", alpha=0.7, edgecolor="none", pad=1),
        )

        # 5. Plot Mission Event POIs
        for poi in mission_event_pois:
            # No proximity filtering for mission events - they are critical to show
            # even if they overlap (though rare).

            # Plot waypoint marker (Blue Diamond)
            ax.plot(
                poi["lon"],
                poi["lat"],
                marker="D",
                color="#3498db",
                markersize=8,
                markeredgecolor="white",
                markeredgewidth=1.5,
                transform=ccrs.PlateCarree(),
                zorder=10,
            )

            # Label with waypoint name
            if poi["label"]:
                ax.text(
                    poi["lon"] + 0.5,
                    poi["lat"] + 0.5,
                    poi["label"],
                    transform=ccrs.PlateCarree(),
                    fontsize=8,
                    fontweight="bold",
                    color="#2c3e50",
                    zorder=11,
                    bbox=dict(facecolor="white", alpha=0.6, edgecolor="none", pad=0.5),
                )

    # Phase 12: Add legend inset to map
    legend_elements = [
        # Route status colors
        Line2D([0], [0], color=STATUS_COLORS["nominal"], linewidth=3, label="Nominal"),
        Line2D(
            [0], [0], color=STATUS_COLORS["degraded"], linewidth=3, label="Degraded"
        ),
        Line2D(
            [0], [0], color=STATUS_COLORS["critical"], linewidth=3, label="Critical"
        ),
        # Marker types
        Line2D(
            [0],
            [0],
            marker="o",
            color="w",
            markerfacecolor="#2ecc71",
            markersize=8,
            label="Departure",
            linestyle="None",
        ),
        Line2D(
            [0],
            [0],
            marker="o",
            color="w",
            markerfacecolor="#e74c3c",
            markersize=8,
            label="Arrival",
            linestyle="None",
        ),
        Line2D(
            [0],
            [0],
            marker="D",
            color="w",
            markerfacecolor="#3498db",
            markersize=8,
            label="POI",
            linestyle="None",
        ),
    ]

    # Add legend inset at lower-right, positioned to not extend beyond figure
    legend = ax.legend(
        handles=legend_elements,
        loc="lower right",
        fontsize=9,
        framealpha=0.95,
        edgecolor="#2c3e50",
        fancybox=True,
    )
    # Ensure legend is drawn on top and doesn't extend beyond axes
    legend.set_zorder(100)

    # Save to PNG bytes
    buf = io.BytesIO()
    plt.savefig(buf, format="png", dpi=dpi, bbox_inches="tight", pad_inches=0)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _generate_timeline_chart(timeline: MissionLegTimeline) -> bytes:
    """Generate a PNG image of a horizontal bar chart showing transport timeline.

    Args:
        timeline: The mission timeline with segments containing transport states

    Returns:
        PNG image as bytes showing three rows of colored blocks representing
        X-Band, Ka (CommKa), and Ku (StarShield) state transitions over mission duration.
    """
    from matplotlib.ticker import FuncFormatter

    # Handle empty timeline
    if not timeline.segments:
        fig, ax = plt.subplots(figsize=(16, 5), dpi=200, facecolor="white")
        ax.text(
            0.5,
            0.5,
            "No timeline segments available",
            ha="center",
            va="center",
            fontsize=14,
            transform=ax.transAxes,
            color="#666",
        )
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        ax.set_facecolor("#f8f9fa")
        ax.axis("off")

        buffer = io.BytesIO()
        fig.savefig(
            buffer, format="png", dpi=200, bbox_inches="tight", facecolor="white"
        )
        plt.close(fig)
        buffer.seek(0)
        return buffer.read()

    # Get mission duration from last segment end time
    mission_end_seconds = (
        ensure_timezone(timeline.segments[-1].end_time)
        - ensure_timezone(timeline.segments[0].start_time)
    ).total_seconds()

    # Create figure with better aspect ratio (16:5 instead of 14:4)
    fig, ax = plt.subplots(figsize=(16, 5), dpi=200, facecolor="white")

    # Modern state to color mapping
    state_colors = {
        "AVAILABLE": "#27ae60",  # Modern green
        "DEGRADED": "#f39c12",  # Modern amber
        "OFFLINE": "#e74c3c",  # Modern red
    }

    # Transport configuration: (y_position, name, state_getter)
    transports = [
        (0, "Ku (StarShield)", lambda seg: seg.ku_state),
        (1, "Ka (CommKa)", lambda seg: seg.ka_state),
        (2, "X-Band", lambda seg: seg.x_state),
    ]

    mission_start = ensure_timezone(timeline.segments[0].start_time)

    # Draw segments for each transport
    for y_pos, transport_name, state_getter in transports:
        for segment in timeline.segments:
            seg_start = ensure_timezone(segment.start_time)
            seg_end = ensure_timezone(segment.end_time)

            # Calculate position and width in mission seconds
            start_offset = (seg_start - mission_start).total_seconds()
            duration = (seg_end - seg_start).total_seconds()

            # Get state and color
            state = state_getter(segment)
            state_str = (
                state.value.upper() if hasattr(state, "value") else str(state).upper()
            )
            color = state_colors.get(state_str, "#808080")

            # Draw horizontal bar with modern styling
            ax.barh(
                y_pos,
                duration,
                left=start_offset,
                height=0.7,
                color=color,
                edgecolor="#2c3e50",
                linewidth=1,
                alpha=0.95,
            )

    # Configure y-axis with better spacing and readability
    ax.set_ylim(-0.6, 2.6)
    ax.set_yticks([0, 1, 2])
    ax.set_yticklabels(
        ["Ku (StarShield)", "Ka (CommKa)", "X-Band"],
        fontsize=11,
        fontweight="semibold",
        color="#2c3e50",
    )
    ax.tick_params(axis="y", labelsize=11, colors="#2c3e50")

    # Configure x-axis with T+ formatting
    def format_time_label(seconds, pos):
        hours = int(seconds // 3600)
        minutes = int((seconds % 3600) // 60)
        return f"T+{hours:02d}:{minutes:02d}"

    ax.xaxis.set_major_formatter(FuncFormatter(format_time_label))
    ax.set_xlim(0, mission_end_seconds)

    # Add vertical grid lines at 1-hour intervals with modern styling
    for hour in range(0, int(mission_end_seconds) + 1, 3600):
        if hour <= mission_end_seconds:
            ax.axvline(x=hour, color="#bdc3c7", linestyle="-", linewidth=0.5, alpha=0.4)

    # Configure grid and labels with modern styling
    ax.grid(True, axis="x", alpha=0.2, linestyle="-", color="#bdc3c7")
    ax.set_axisbelow(True)
    ax.set_xlabel("Mission Time", fontsize=12, fontweight="bold", color="#2c3e50")
    ax.set_title(
        "Transport State Timeline",
        fontsize=14,
        fontweight="bold",
        pad=20,
        color="#2c3e50",
    )

    # Style x-axis
    ax.tick_params(axis="x", labelsize=10, colors="#2c3e50")

    # Add modern legend
    legend_elements = [
        Line2D([0], [0], color="#27ae60", lw=12, label="Available", alpha=0.95),
        Line2D([0], [0], color="#f39c12", lw=12, label="Degraded", alpha=0.95),
        Line2D([0], [0], color="#e74c3c", lw=12, label="Offline", alpha=0.95),
    ]
    ax.legend(
        handles=legend_elements,
        loc="upper right",
        fontsize=10,
        framealpha=0.95,
        edgecolor="#bdc3c7",
        fancybox=True,
    )

    # Style the plot background
    ax.set_facecolor("#f8f9fa")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#bdc3c7")
    ax.spines["bottom"].set_color("#bdc3c7")

    # Adjust layout
    plt.tight_layout()

    # Save to buffer with high quality
    buffer = io.BytesIO()
    fig.savefig(buffer, format="png", dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buffer.seek(0)
    return buffer.read()


def _segment_rows(
    timeline: MissionLegTimeline,
    mission: Mission | None,
) -> pd.DataFrame:
    """Convert timeline segments into a pandas DataFrame."""
    mission_start = mission_start_timestamp(timeline)
    rows: list[tuple[datetime, int, dict]] = []

    for idx, segment in enumerate(timeline.segments, start=1):
        start_utc = ensure_timezone(segment.start_time)
        end_value = segment.end_time if segment.end_time else segment.start_time
        end_utc = ensure_timezone(end_value)
        duration_seconds = max((end_utc - start_utc).total_seconds(), 0)
        warning_only = segment_is_x_ku_warning(segment)
        status_value = (
            segment.status.value
            if isinstance(segment.status, TimelineStatus)
            else str(segment.status)
        )
        status_value = status_value.upper()
        impacted_display = serialize_transport_list(segment.impacted_transports)
        if warning_only:
            status_value = TimelineStatus.NOMINAL.value.upper()
            impacted_display = ""
        record = {
            "Segment #": idx,
            "Mission ID": mission.id if mission else timeline.mission_id,
            "Mission Name": (
                mission.name if mission and mission.name else timeline.mission_id
            ),
            "Status": status_value,
            "Start Time": compose_time_block(start_utc, mission_start),
            "End Time": compose_time_block(end_utc, mission_start),
            "Duration": format_seconds_hms(duration_seconds),
            TRANSPORT_DISPLAY[Transport.X]: display_transport_state(
                segment.x_state, warning_override=warning_only
            ),
            TRANSPORT_DISPLAY[Transport.KA]: segment.ka_state.value.upper(),
            TRANSPORT_DISPLAY[Transport.KU]: segment.ku_state.value.upper(),
            "Impacted Transports": impacted_display,
            "Reasons": ", ".join(segment.reasons),
            "Metadata": (
                json.dumps(segment.metadata, sort_keys=True) if segment.metadata else ""
            ),
        }
        rows.append((start_utc, 1, record))

    rows.extend(_aar_block_rows(timeline, mission, mission_start))

    columns = [
        "Segment #",
        "Mission ID",
        "Mission Name",
        "Status",
        "Start Time",
        "End Time",
        "Duration",
        TRANSPORT_DISPLAY[Transport.X],
        TRANSPORT_DISPLAY[Transport.KA],
        TRANSPORT_DISPLAY[Transport.KU],
        "Impacted Transports",
        "Reasons",
        "Metadata",
    ]

    ordered_records = [
        row for _, _, row in sorted(rows, key=lambda item: (item[0], item[1]))
    ]
    return pd.DataFrame.from_records(ordered_records, columns=columns)


def _advisory_rows(
    timeline: MissionLegTimeline, mission: Mission | None
) -> pd.DataFrame:
    """Convert timeline advisories into a pandas DataFrame (may be empty)."""
    mission_start = mission_start_timestamp(timeline)
    records: list[dict] = []
    for advisory in timeline.advisories:
        ts_utc = ensure_timezone(advisory.timestamp)
        records.append(
            {
                "Mission ID": mission.id if mission else timeline.mission_id,
                "Timestamp (UTC)": format_utc(ts_utc),
                "Timestamp (Eastern)": format_eastern(ts_utc),
                "T Offset": format_offset(ts_utc - mission_start),
                "Transport": (
                    advisory.transport.value
                    if isinstance(advisory.transport, Transport)
                    else advisory.transport
                ),
                "Severity": advisory.severity,
                "Event Type": advisory.event_type,
                "Message": advisory.message,
                "Metadata": (
                    json.dumps(advisory.metadata, sort_keys=True)
                    if advisory.metadata
                    else ""
                ),
            }
        )
    columns = [
        "Mission ID",
        "Timestamp (UTC)",
        "Timestamp (Eastern)",
        "T Offset",
        "Transport",
        "Severity",
        "Event Type",
        "Message",
        "Metadata",
    ]
    return pd.DataFrame.from_records(records, columns=columns)


def _statistics_rows(timeline: MissionLegTimeline) -> pd.DataFrame:
    stats = timeline.statistics or {}
    if not stats:
        return pd.DataFrame(columns=["Metric", "Value"])
    rows = []
    for key, value in stats.items():
        if key.startswith("_"):
            continue
        display_name = humanize_metric_name(key)
        display_value = value
        if isinstance(value, (int, float)) and key.endswith("seconds"):
            display_value = format_seconds_hms(value)
        rows.append({"Metric": display_name, "Value": display_value})
    return pd.DataFrame(rows, columns=["Metric", "Value"])


# Note: format_seconds_hms, humanize_metric_name, and compose_time_block are now
# imported from exporter.formatting module above


def generate_csv_export(
    timeline: MissionLegTimeline, mission: Mission | None = None
) -> bytes:
    """Return CSV bytes for the mission timeline."""
    csv_buffer = io.StringIO()
    df = _segment_rows(timeline, mission)
    df.to_csv(csv_buffer, index=False)
    return csv_buffer.getvalue().encode("utf-8")


def _summary_table_rows(
    timeline: MissionLegTimeline, mission: Mission | None = None
) -> pd.DataFrame:
    """Generate simplified summary table DataFrame with key columns for Sheet 1.

    Returns a DataFrame with columns: Start (UTC), Duration, Status, Systems Down
    """
    records: list[dict] = []

    for segment in timeline.segments:
        start_utc = ensure_timezone(segment.start_time)
        end_value = segment.end_time if segment.end_time else segment.start_time
        end_utc = ensure_timezone(end_value)
        duration_seconds = max((end_utc - start_utc).total_seconds(), 0)

        # Get segment status
        status_value = (
            segment.status.value
            if isinstance(segment.status, TimelineStatus)
            else str(segment.status)
        )
        status_value = status_value.upper()

        # Get systems down (impacted transports)
        systems_down = serialize_transport_list(segment.impacted_transports)

        record = {
            "Start (UTC)": format_utc(start_utc),
            "Duration": format_seconds_hms(duration_seconds),
            "Status": status_value,
            "Systems Down": systems_down,
        }
        records.append(record)

    columns = ["Start (UTC)", "Duration", "Status", "Systems Down"]
    return pd.DataFrame.from_records(records, columns=columns)


def generate_xlsx_export(
    timeline: MissionLegTimeline,
    mission: Mission | None = None,
    parent_mission_id: str | None = None,
    route_manager: RouteManager | None = None,
    poi_manager: POIManager | None = None,
) -> bytes:
    """Return XLSX bytes containing Summary sheet (with map, chart, table) plus Timeline, Advisory, and Statistics sheets."""
    workbook_bytes = io.BytesIO()

    # Generate all data
    summary_df = _summary_table_rows(timeline, mission)
    timeline_df = _segment_rows(timeline, mission)
    advisories_df = _advisory_rows(timeline, mission)
    stats_df = _statistics_rows(timeline)
    map_image_bytes = _generate_route_map(
        timeline,
        mission,
        parent_mission_id=parent_mission_id,
        route_manager=route_manager,
        poi_manager=poi_manager,
    )
    chart_image_bytes = _generate_timeline_chart(timeline)

    # Write DataFrames to Excel sheets
    with pd.ExcelWriter(workbook_bytes, engine="openpyxl") as writer:
        # Write Summary sheet first (will reorder later)
        summary_df.to_excel(writer, sheet_name="Summary", index=False, startrow=0)
        timeline_df.to_excel(writer, sheet_name="Timeline", index=False)
        if not advisories_df.empty:
            advisories_df.to_excel(writer, sheet_name="Advisories", index=False)
        if not stats_df.empty:
            stats_df.to_excel(writer, sheet_name="Statistics", index=False)

        # Access openpyxl workbook for image embedding and formatting
        ws_summary = writer.sheets["Summary"]

        # Set column widths for summary sheet
        ws_summary.column_dimensions["A"].width = 25
        ws_summary.column_dimensions["B"].width = 15
        ws_summary.column_dimensions["C"].width = 12
        ws_summary.column_dimensions["D"].width = 30

        # Insert rows at top to accommodate map and chart images
        # Map will be at A1 (approximately 30 rows tall)
        # Chart will be at A32 (approximately 15 rows tall)
        # Table will be at A49
        ws_summary.insert_rows(1, 48)

        # Embed map image at A1
        try:
            map_image_stream = io.BytesIO(map_image_bytes)
            map_image = OpenpyxlImage(map_image_stream)
            map_image.width = 750  # pixels
            map_image.height = 500  # pixels
            ws_summary.add_image(map_image, "A1")
        except Exception as e:
            logger.error("Failed to embed map image in Excel: %s", e, exc_info=True)

        # Embed timeline chart at A32
        try:
            chart_image_stream = io.BytesIO(chart_image_bytes)
            chart_image = OpenpyxlImage(chart_image_stream)
            chart_image.width = 850  # pixels
            chart_image.height = 300  # pixels
            ws_summary.add_image(chart_image, "A32")
        except Exception as e:
            logger.error(
                "Failed to embed timeline chart in Excel: %s", e, exc_info=True
            )

        # Apply color formatting to summary table rows (starting at row 49)
        # Header is at row 49, data starts at row 50
        color_map = {
            "NOMINAL": PatternFill(
                start_color="00FF00", end_color="00FF00", fill_type="solid"
            ),
            "DEGRADED": PatternFill(
                start_color="FFFF00", end_color="FFFF00", fill_type="solid"
            ),
            "CRITICAL": PatternFill(
                start_color="FF0000", end_color="FF0000", fill_type="solid"
            ),
        }

        # Apply color to data rows (skip header row)
        for row_idx, (_, row) in enumerate(summary_df.iterrows(), start=50):
            status = row["Status"]
            fill = color_map.get(status)
            if fill:
                # Apply fill to columns A-D for this row
                for col_idx in range(1, 5):  # Columns A-D (1-4)
                    ws_summary.cell(row=row_idx, column=col_idx).fill = fill

    workbook_bytes.seek(0)
    wb_final = openpyxl.load_workbook(workbook_bytes)

    # Reorder sheets: move Summary to position 0 (first)
    if "Summary" in wb_final.sheetnames:
        summary_sheet = wb_final["Summary"]
        wb_final._sheets.remove(summary_sheet)
        wb_final._sheets.insert(0, summary_sheet)

    # Write final workbook to bytes
    final_bytes = io.BytesIO()
    wb_final.save(final_bytes)
    final_bytes.seek(0)
    return final_bytes.read()


def generate_pdf_export(
    timeline: MissionLegTimeline,
    mission: Mission | None = None,
    parent_mission_id: str | None = None,
    route_manager: RouteManager | None = None,
    poi_manager: POIManager | None = None,
) -> bytes:
    """Render a PDF brief summarizing the mission timeline."""
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=landscape(LETTER),
        leftMargin=0.75 * inch,
        rightMargin=0.75 * inch,
        topMargin=0.75 * inch,
        bottomMargin=0.75 * inch,
    )
    styles = getSampleStyleSheet()
    story: list = []

    mission_name = mission.name if mission and mission.name else timeline.mission_id
    logo_flow = _load_logo_flowable()
    if logo_flow:
        logo_width = logo_flow.drawWidth
        header_table = Table(
            [
                [
                    Paragraph("Mission Communication Timeline", styles["Title"]),
                    logo_flow,
                ],
                [
                    Paragraph(f"Mission: {mission_name}", styles["Heading2"]),
                    Spacer(1, 0),
                ],
            ],
            colWidths=[doc.width - logo_width, logo_width],
        )
        header_table.setStyle(
            TableStyle(
                [
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("ALIGN", (1, 0), (1, 0), "RIGHT"),
                    ("SPAN", (1, 1), (1, 1)),
                ]
            )
        )
        story.append(header_table)
    else:
        story.append(Paragraph("Mission Communication Timeline", styles["Title"]))
        story.append(Spacer(1, 0.05 * inch))
        story.append(Paragraph(f"Mission: {mission_name}", styles["Heading2"]))
    story.append(Spacer(1, 0.2 * inch))

    stats_df = _statistics_rows(timeline)
    if not stats_df.empty:
        stats_table = Table(
            [["Metric", "Value"], *stats_df.values.tolist()],
            hAlign="LEFT",
        )
        stats_table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.black),
                    ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("BOTTOMPADDING", (0, 0), (-1, 0), 6),
                ]
            )
        )
        story.append(stats_table)
        story.append(Spacer(1, 0.2 * inch))

    story.append(Spacer(1, 0.3 * inch))

    # Add route map
    story.append(Paragraph("Route Map", styles["Heading2"]))
    try:
        route_map_bytes = _generate_route_map(
            timeline,
            mission,
            parent_mission_id=parent_mission_id,
            route_manager=route_manager,
            poi_manager=poi_manager,
        )
        route_map_stream = io.BytesIO(route_map_bytes)
        route_map_image = Image(route_map_stream, width=6.5 * inch, height=4.3 * inch)
        story.append(route_map_image)
    except Exception as e:
        logger.error("Failed to generate route map for PDF: %s", e, exc_info=True)
        story.append(Paragraph(f"[Route map unavailable: {str(e)}]", styles["Normal"]))

    # Page break before timeline
    story.append(PageBreak())

    # Add timeline chart
    story.append(Paragraph("Transport Timeline", styles["Heading2"]))
    try:
        chart_image_bytes = _generate_timeline_chart(timeline)
        chart_image_stream = io.BytesIO(chart_image_bytes)
        chart_image = Image(chart_image_stream, width=7 * inch, height=2.3 * inch)
        story.append(chart_image)
    except Exception as e:
        logger.error("Failed to generate timeline chart for PDF: %s", e, exc_info=True)
        story.append(
            Paragraph(f"[Timeline chart unavailable: {str(e)}]", styles["Normal"])
        )

    story.append(Spacer(1, 0.2 * inch))

    timeline_df = _segment_rows(timeline, mission)
    if timeline_df.empty:
        story.append(
            Paragraph("No timeline segments were generated.", styles["Normal"])
        )
    else:
        story.append(Paragraph("Timeline Segments", styles["Heading2"]))
        table_columns = [
            ("", "Segment #"),
            ("Status", "Status"),
            ("Start Time", "Start Time"),
            ("End Time", "End Time"),
            ("Duration", "Duration"),
            (TRANSPORT_DISPLAY[Transport.X], TRANSPORT_DISPLAY[Transport.X]),
            (TRANSPORT_DISPLAY[Transport.KA], TRANSPORT_DISPLAY[Transport.KA]),
            (TRANSPORT_DISPLAY[Transport.KU], TRANSPORT_DISPLAY[Transport.KU]),
            ("Reasons", "Reasons"),
        ]
        body_style = styles["BodyText"]
        body_style.fontSize = 8
        body_style.leading = 9.5
        # Pre-calculate status overrides and identify critical rows
        # We need to modify the data *before* creating the Table if we want to change text

        # Create a map of row_idx -> is_critical for styling later
        critical_rows = {}

        # Rebuild data with overrides
        data = [[header for header, _ in table_columns]]
        for row_idx, row in timeline_df.iterrows():
            # Determine if row is critical (2+ bad systems, excluding SoF)
            bad_cols_count = 0
            for name in STATE_COLUMNS:
                val = str(row[name]).strip().lower()
                if val in ("degraded", "warning", "offline"):
                    bad_cols_count += 1

            is_critical = bad_cols_count >= 2
            critical_rows[row_idx] = is_critical

            # Check for SoF/AAR reasons
            reasons = str(row.get("Reasons", "")).lower()
            is_sof_row = "safety-of-flight" in reasons or "aar" in reasons

            row_values = []
            for header, key in table_columns:
                value = row[key]

                # Override Status text
                if key == "Status":
                    if is_critical:
                        value = "CRITICAL"
                    elif is_sof_row and str(value).lower() in (
                        "nominal",
                        "available",
                        "warning",
                    ):
                        value = "SOF"

                if key in {"Start Time", "End Time"}:
                    display = (value or "-").replace("\n", "<br/>")
                    row_values.append(Paragraph(display, body_style))
                elif key == "Reasons":
                    display = value or "-"
                    row_values.append(Paragraph(str(display), body_style))
                else:
                    row_values.append(value)
            data.append(row_values)

        column_weights = [0.5, 1.1, 1.5, 1.5, 0.8, 1.2, 1.2, 1.2, 3.0]
        width_unit = doc.width / sum(column_weights)
        col_widths = [weight * width_unit for weight in column_weights]
        table = Table(data, repeatRows=1, colWidths=col_widths)
        style_commands = [
            ("BACKGROUND", (0, 0), (-1, 0), colors.Color(0.2, 0.4, 0.7)),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, 0), 9),
            ("FONTSIZE", (0, 1), (-1, -1), 8),
            ("ALIGN", (0, 0), (-1, -1), "LEFT"),
            ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.whitesmoke, colors.lightgrey]),
        ]
        column_index = {key: idx for idx, (_, key) in enumerate(table_columns)}
        state_column_indices = {name: column_index[name] for name in STATE_COLUMNS}
        status_col_idx = column_index.get("Status")

        for row_idx in range(len(timeline_df)):
            is_critical = critical_rows.get(row_idx, False)

            # Check for SoF/AAR reasons
            reasons = str(timeline_df.iloc[row_idx]["Reasons"]).lower()
            is_sof_row = "safety-of-flight" in reasons or "aar" in reasons

            # Color Status column
            if status_col_idx is not None:
                status_val = str(timeline_df.iloc[row_idx]["Status"]).upper()
                if is_critical:
                    style_commands.append(
                        (
                            "BACKGROUND",
                            (status_col_idx, row_idx + 1),
                            (status_col_idx, row_idx + 1),
                            LIGHT_RED,
                        )
                    )
                elif status_val in ("DEGRADED", "WARNING") or is_sof_row:
                    style_commands.append(
                        (
                            "BACKGROUND",
                            (status_col_idx, row_idx + 1),
                            (status_col_idx, row_idx + 1),
                            LIGHT_YELLOW,
                        )
                    )

            # Color Transport columns
            for name in STATE_COLUMNS:
                val = str(timeline_df.iloc[row_idx][name]).strip().lower()
                col_idx = state_column_indices[name]

                # Check if this specific cell is SoF (AVAILABLE but needs yellow)
                # SoF applies if the row is SoF and this cell is AVAILABLE/NOMINAL
                cell_is_sof = is_sof_row and val in ("available", "nominal")

                if val in ("degraded", "warning", "offline") or cell_is_sof:
                    if is_critical and not cell_is_sof:
                        color = LIGHT_RED
                    elif val in ("degraded", "warning") or cell_is_sof:
                        color = LIGHT_YELLOW
                    else:
                        color = LIGHT_RED  # Offline

                    style_commands.append(
                        (
                            "BACKGROUND",
                            (col_idx, row_idx + 1),
                            (col_idx, row_idx + 1),
                            color,
                        )
                    )

        table.setStyle(TableStyle(style_commands))
        story.append(table)

    footer_style = styles["Normal"].clone("Footer")
    footer_style.alignment = TA_RIGHT
    story.append(Spacer(1, 0.2 * inch))
    story.append(
        Paragraph(
            f"Timeline generated: {format_utc(timeline.created_at)}", footer_style
        )
    )

    doc.build(story)
    buffer.seek(0)
    return buffer.read()


def generate_pptx_export(
    timeline: MissionLegTimeline,
    mission: Mission | None = None,
    parent_mission_id: str | None = None,
    route_manager: RouteManager | None = None,
    poi_manager: POIManager | None = None,
) -> bytes:
    """Generate a PowerPoint presentation with map and timeline table."""
    prs = Presentation()

    # Slide 1: Route Map
    # Use a blank layout (usually index 6 in default template)
    blank_slide_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_slide_layout)

    # Generate map image
    try:
        map_image_bytes = _generate_route_map(
            timeline,
            mission,
            parent_mission_id=parent_mission_id,
            route_manager=route_manager,
            poi_manager=poi_manager,
        )
        map_image_stream = io.BytesIO(map_image_bytes)

        # Add image to slide, scaled to fit
        # Default slide size is 10x7.5 inches
        # We want to maximize the map visibility
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(6.5)

        slide1.shapes.add_picture(
            map_image_stream, left, top, width=width, height=height
        )

        # Add title
        title_box = slide1.shapes.add_textbox(
            Inches(0.5), Inches(0.1), Inches(9), Inches(0.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Mission Route Map"
        p.font.bold = True
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

    except Exception as e:
        logger.error("Failed to add map to PPTX: %s", e, exc_info=True)
        textbox = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        textbox.text = f"Map generation failed: {str(e)}"

    # Slide 2: Timeline Table
    # Slide 2+: Timeline Table (Paginated)
    timeline_df = _segment_rows(timeline, mission)

    if not timeline_df.empty:
        # Define columns to show
        columns_to_show = [
            "Segment #",
            "Status",
            "Start Time",
            "End Time",
            "Duration",
            TRANSPORT_DISPLAY[Transport.X],
            TRANSPORT_DISPLAY[Transport.KA],
            TRANSPORT_DISPLAY[Transport.KU],
            "Reasons",
        ]

        # Pagination settings
        ROWS_PER_SLIDE = 10
        MIN_ROWS_LAST_SLIDE = 3

        # Convert to list of dicts for easier manipulation if needed, or just slice dataframe
        # Slicing dataframe is easier but we need to handle the orphan logic
        records = timeline_df.to_dict("records")
        total_rows = len(records)

        chunks = []
        if total_rows <= ROWS_PER_SLIDE:
            chunks.append(timeline_df)
        else:
            # Naive split first
            current_idx = 0
            while current_idx < total_rows:
                end_idx = min(current_idx + ROWS_PER_SLIDE, total_rows)
                chunks.append(timeline_df.iloc[current_idx:end_idx])
                current_idx = end_idx

            # Check last chunk
            if len(chunks) > 1:
                last_chunk = chunks[-1]
                if len(last_chunk) < MIN_ROWS_LAST_SLIDE:
                    # Need to move items from second to last chunk
                    # Re-slice the last two chunks
                    # We need to go back to the source dataframe indices
                    # Let's recalculate the split point for the last page

                    # Total items excluding the last "naive" chunk's items
                    # Actually, let's just do it numerically
                    # If we have 11 items. 10, 1.
                    # We want 8, 3.
                    # Split point should be total - 3.

                    # If we have 21 items. 10, 10, 1.
                    # We want 10, 8, 3.

                    # So we only adjust the boundary between the last two chunks.
                    # The start of the last chunk is currently `total_rows - len(last_chunk)`
                    # We want the start to be `total_rows - 3` (or whatever min is)
                    # But we must ensure the second to last chunk doesn't exceed max? No, it will shrink.
                    # We must ensure the second to last chunk doesn't fall below min?
                    # If max=10, min=3.
                    # Worst case for prev chunk: It had 10. We take 2. It has 8. Safe.

                    # Re-build chunks
                    # Everything up to the second to last chunk stays same
                    base_chunks = chunks[:-2]

                    # Combine last two
                    remainder_count = len(chunks[-2]) + len(chunks[-1])
                    # We want last chunk to have MIN_ROWS_LAST_SLIDE
                    # So second to last has remainder_count - MIN_ROWS_LAST_SLIDE

                    split_idx = (
                        total_rows - remainder_count
                    )  # Start of the second-to-last chunk
                    second_last_len = remainder_count - MIN_ROWS_LAST_SLIDE

                    chunk_second_last = timeline_df.iloc[
                        split_idx : split_idx + second_last_len
                    ]
                    chunk_last = timeline_df.iloc[split_idx + second_last_len :]

                    chunks = base_chunks + [chunk_second_last, chunk_last]

        for chunk_idx, chunk in enumerate(chunks):
            # Add new slide for each chunk
            slide = prs.slides.add_slide(blank_slide_layout)

            # Add title
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.1), Inches(9), Inches(0.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = (
                "Mission Timeline" if chunk_idx == 0 else "Mission Timeline (cont.)"
            )
            p.font.bold = True
            p.font.size = Pt(24)
            p.alignment = PP_ALIGN.CENTER

            rows = len(chunk) + 1  # +1 for header
            cols = len(columns_to_show)

            # Create table
            left = Inches(0.5)
            top = Inches(1.0)
            width = Inches(9.0)
            # Use a minimal height so rows don't stretch to fill a large area
            # The table will expand as needed
            height = Inches(1.0)

            shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = shape.table

            # Set column widths (adjusted for wider times)
            # Total width 9 inches
            # Old weights: [0.6, 1.0, 1.5, 1.5, 1.0, 1.0, 1.0, 1.0, 2.0]
            # New weights: [0.6, 1.0, 1.75, 1.75, 1.0, 1.0, 1.0, 1.0, 1.5]
            col_weights = [0.6, 1.0, 1.75, 1.75, 1.0, 1.0, 1.0, 1.0, 1.5]
            total_weight = sum(col_weights)
            for i, weight in enumerate(col_weights):
                table.columns[i].width = int(width * (weight / total_weight))

            # Header row style
            for i, col_name in enumerate(columns_to_show):
                cell = table.cell(0, i)
                cell.text = col_name
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(51, 102, 178)  # Dark Blue
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White
                paragraph.font.bold = True
                paragraph.font.size = Pt(10)
                paragraph.alignment = PP_ALIGN.CENTER

            # Data rows
            for row_idx, (_, row_data) in enumerate(chunk.iterrows(), start=1):
                # Pre-calculate bad transports for this row
                bad_transports = []
                for col_name in STATE_COLUMNS:
                    val = str(row_data[col_name] if row_data[col_name] else "").lower()
                    if val in ("degraded", "warning", "offline"):
                        bad_transports.append(col_name)

                is_critical_row = len(bad_transports) >= 2

                for col_idx, col_name in enumerate(columns_to_show):
                    cell = table.cell(row_idx, col_idx)
                    val = str(
                        row_data[col_name] if row_data[col_name] is not None else ""
                    )

                    # Override Status text if critical
                    if col_name == "Status" and is_critical_row:
                        val = "CRITICAL"

                    cell.text = val

                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(9)
                        paragraph.alignment = PP_ALIGN.LEFT

                    # Alternating row colors
                    cell.fill.solid()
                    if row_idx % 2 == 0:
                        cell.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light Gray
                    else:
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White

                    # Status coloring for transport columns and Status column
                    # Status coloring for transport columns and Status column
                    if col_name in STATE_COLUMNS or col_name == "Status":
                        val_lower = val.lower()

                        # Check for Safety-of-Flight in reasons if available/nominal/warning
                        is_sof = False
                        # We check reasons regardless of status, but only apply SOF override if appropriate
                        reasons = str(row_data.get("Reasons", "")).lower()
                        if "safety-of-flight" in reasons or "aar" in reasons:
                            is_sof = True

                        # Override Status text to "SOF" if it's a safety window
                        # User request: "baseline for both should just be SOF"
                        # If status is WARNING (e.g. due to AAR) or NOMINAL/AVAILABLE, show SOF.
                        # If DEGRADED or OFFLINE, show that.
                        if (
                            col_name == "Status"
                            and is_sof
                            and val_lower in ("available", "nominal", "warning")
                        ):
                            cell.text = "SOF"
                            # Re-apply font size since setting text resets it
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.font.size = Pt(9)
                                paragraph.alignment = PP_ALIGN.LEFT

                        if val_lower in ("degraded", "warning", "offline") or is_sof:
                            cell.fill.solid()
                            if is_critical_row and not is_sof:
                                cell.fill.fore_color.rgb = RGBColor(
                                    255, 204, 204
                                )  # Light Red
                            elif val_lower in ("degraded", "warning") or is_sof:
                                cell.fill.fore_color.rgb = RGBColor(
                                    255, 255, 204
                                )  # Light Yellow
                            else:
                                cell.fill.fore_color.rgb = RGBColor(
                                    255, 204, 204
                                )  # Light Red

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()


def _base_map_canvas() -> bytes:
    """Generate a blank 4K map canvas with no route or markers."""
    width_inches = 3840 / 300
    height_inches = 2880 / 300

    fig = plt.figure(figsize=(width_inches, height_inches), dpi=300, facecolor="white")
    ax = fig.add_subplot(111, projection=ccrs.PlateCarree())
    fig.subplots_adjust(left=0, right=1, top=1, bottom=0)
    ax.set_extent([-180, 180, -90, 90], crs=ccrs.PlateCarree())
    ax.coastlines(resolution="50m", linewidth=0.5, color="#2c3e50")
    ax.add_feature(cfeature.BORDERS, linewidth=0.5, color="#bdc3c7")
    ax.add_feature(cfeature.LAND, facecolor="#ecf0f1", edgecolor="none")
    ax.add_feature(cfeature.OCEAN, facecolor="#d5e8f7", edgecolor="none")
    ax.gridlines(
        draw_labels=False, alpha=0.1, linestyle="-", linewidth=0.3, color="#95a5a6"
    )
    ax.spines["geo"].set_visible(False)
    ax.set_xticks([])
    ax.set_yticks([])

    buf = io.BytesIO()
    plt.savefig(buf, format="png", dpi=300, bbox_inches="tight", pad_inches=0)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def generate_timeline_export(
    export_format: TimelineExportFormat,
    mission: Mission,
    timeline: MissionLegTimeline,
    parent_mission_id: str | None = None,
    route_manager: RouteManager | None = None,
    poi_manager: POIManager | None = None,
) -> ExportArtifact:
    """Generate the requested export artifact."""
    if export_format is TimelineExportFormat.CSV:
        content = generate_csv_export(timeline, mission)
        return ExportArtifact(content=content, media_type="text/csv", extension="csv")
    if export_format is TimelineExportFormat.XLSX:
        content = generate_xlsx_export(
            timeline,
            mission,
            parent_mission_id=parent_mission_id,
            route_manager=route_manager,
            poi_manager=poi_manager,
        )
        media_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        return ExportArtifact(content=content, media_type=media_type, extension="xlsx")
    if export_format is TimelineExportFormat.PDF:
        content = generate_pdf_export(
            timeline,
            mission,
            parent_mission_id=parent_mission_id,
            route_manager=route_manager,
            poi_manager=poi_manager,
        )
        return ExportArtifact(
            content=content, media_type="application/pdf", extension="pdf"
        )
    if export_format is TimelineExportFormat.PPTX:
        content = generate_pptx_export(
            timeline,
            mission,
            parent_mission_id=parent_mission_id,
            route_manager=route_manager,
            poi_manager=poi_manager,
        )
        return ExportArtifact(
            content=content,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            extension="pptx",
        )

    raise ExportGenerationError(f"Unsupported export format: {export_format}")
