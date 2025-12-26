"""PowerPoint styling utilities for mission export slides.

This module provides reusable styling functions to create professional, branded
PowerPoint presentations with consistent colors, layouts, and visual hierarchy.

Color palette based on organizational branding standards:
- Gold (RGB 212, 175, 55): Primary accent for headers and branding
- Light gray (RGB 248, 249, 250): Content backgrounds
- Status colors: Green (NOMINAL), Blue (SOF), Orange (DEGRADED), Red (CRITICAL)

All functions accept `slide` objects from python-pptx and return created shapes
or boolean success indicators.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import TYPE_CHECKING

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

if TYPE_CHECKING:
    from pptx.slide import Slide
    from app.mission.models import TimelineStatus

logger = logging.getLogger(__name__)

# Brand Colors
BRAND_GOLD = RGBColor(212, 175, 55)  # Primary accent
CONTENT_GRAY = RGBColor(248, 249, 250)  # Light backgrounds

# Status Colors
STATUS_NOMINAL = RGBColor(22, 163, 74)  # Green
STATUS_SOF = RGBColor(2, 132, 199)  # Blue
STATUS_DEGRADED = RGBColor(234, 88, 12)  # Orange
STATUS_CRITICAL = RGBColor(220, 38, 38)  # Red

# Text Colors
TEXT_BLACK = RGBColor(0, 0, 0)
TEXT_WHITE = RGBColor(255, 255, 255)


def add_header_bar(slide: Slide, left: float, top: float, width: float, height: float):
    """Add gold header bar to slide.

    Args:
        slide: PowerPoint slide object
        left: Left position (in Inches)
        top: Top position (in Inches)
        width: Bar width (in Inches)
        height: Bar height (in Inches)

    Returns:
        Shape object representing the header bar
    """
    shape = slide.shapes.add_shape(
        1,  # Rectangle shape type
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_GOLD
    shape.line.fill.background()  # No border
    return shape


def add_footer_bar(slide: Slide, left: float, top: float, width: float, height: float):
    """Add gold footer bar to slide.

    Args:
        slide: PowerPoint slide object
        left: Left position (in Inches)
        top: Top position (in Inches)
        width: Bar width (in Inches)
        height: Bar height (in Inches)

    Returns:
        Shape object representing the footer bar
    """
    shape = slide.shapes.add_shape(
        1,  # Rectangle shape type
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_GOLD
    shape.line.fill.background()  # No border
    return shape


def add_slide_title(slide: Slide, text: str, top: float = 0.25) -> None:
    """Add styled slide title (24pt bold, centered).

    Args:
        slide: PowerPoint slide object
        text: Title text to display
        top: Top position in inches (default: 0.25)
    """
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(top), Inches(9.0), Inches(0.5)
    )
    text_frame = title_box.text_frame
    text_frame.text = text

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(24)
    paragraph.font.bold = True
    paragraph.font.color.rgb = TEXT_BLACK


def add_footer_text(
    slide: Slide,
    text: str,
    bottom: float = 5.5,
    font_size: int = 10,
    color: RGBColor | None = None,
) -> None:
    """Add styled footer text (customizable size and color, centered).

    Args:
        slide: PowerPoint slide object
        text: Footer text to display
        bottom: Bottom position in inches (default: 5.5)
        font_size: Font size in points (default: 10)
        color: Text color as RGBColor (default: TEXT_BLACK)
    """
    footer_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(bottom), Inches(9.0), Inches(0.3)
    )
    text_frame = footer_box.text_frame
    text_frame.text = text

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = color if color else TEXT_BLACK


def add_content_background(
    slide: Slide, left: float, top: float, width: float, height: float
):
    """Add light gray content background.

    Args:
        slide: PowerPoint slide object
        left: Left position (in Inches)
        top: Top position (in Inches)
        width: Background width (in Inches)
        height: Background height (in Inches)

    Returns:
        Shape object representing the content background
    """
    shape = slide.shapes.add_shape(
        1,  # Rectangle shape type
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CONTENT_GRAY
    shape.line.fill.background()  # No border
    return shape


def add_status_badge(
    slide: Slide,
    status: TimelineStatus,
    left: float,
    top: float,
    width: float,
    height: float,
):
    """Add color-coded status badge with white text.

    Args:
        slide: PowerPoint slide object
        status: TimelineStatus enum value
        left: Left position (in Inches)
        top: Top position (in Inches)
        width: Badge width (in Inches)
        height: Badge height (in Inches)

    Returns:
        Shape object representing the status badge
    """
    from app.mission.models import TimelineStatus

    STATUS_COLOR_MAP = {
        TimelineStatus.NOMINAL: STATUS_NOMINAL,
        TimelineStatus.SOF: STATUS_SOF,
        TimelineStatus.DEGRADED: STATUS_DEGRADED,
        TimelineStatus.CRITICAL: STATUS_CRITICAL,
    }

    color = STATUS_COLOR_MAP.get(status, STATUS_NOMINAL)

    shape = slide.shapes.add_shape(
        1,  # Rectangle shape type
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # No border

    # Add status text
    text_frame = shape.text_frame
    text_frame.text = status.value

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(11)
    paragraph.font.bold = True
    paragraph.font.color.rgb = TEXT_WHITE

    return shape


def add_segment_separator(slide: Slide, left: float, top: float, width: float):
    """Add visual separator between timeline segments.

    Args:
        slide: PowerPoint slide object
        left: Left position (in Inches)
        top: Top position (in Inches)
        width: Separator width (in Inches)

    Returns:
        Shape object representing the separator line
    """
    # Use a thin gray rectangle as a separator
    shape = slide.shapes.add_shape(
        1,  # Rectangle shape type
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.01),  # Very thin line
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(200, 200, 200)  # Light gray
    shape.line.fill.background()  # No border
    return shape


def add_logo(
    slide: Slide,
    logo_path: Path,
    left: float,
    top: float,
    max_width: float,
    max_height: float,
) -> bool:
    """Add logo if exists, maintaining aspect ratio.

    Args:
        slide: PowerPoint slide object
        logo_path: Path to logo image file
        left: Left position (in Inches)
        top: Top position (in Inches)
        max_width: Maximum width (in Inches)
        max_height: Maximum height (in Inches)

    Returns:
        True if logo was added, False if logo file not found
    """
    logger.info(f"Attempting to add logo from: {logo_path}")
    logger.info(f"Logo exists: {logo_path.exists()}")

    if not logo_path.exists():
        logger.warning(f"Logo not found at {logo_path}, skipping")
        return False

    try:
        # Add picture with height constraint (maintains aspect ratio)
        # Use the smaller of max_width and max_height to ensure it fits
        size = min(max_width, max_height)
        slide.shapes.add_picture(
            str(logo_path),
            Inches(left),
            Inches(top),
            height=Inches(size),
        )
        logger.info(f"Logo successfully added from {logo_path}")
        return True
    except Exception as e:
        logger.error(f"Failed to add logo: {e}", exc_info=True)
        return False
