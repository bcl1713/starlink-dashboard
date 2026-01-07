"""Mission package export utilities for creating portable mission archives."""

# FR-004: File exceeds 300 lines (1203 lines) because packaging mission data
# involves multi-format export, KML generation, metadata handling, and zip archive
# assembly that are tightly coupled. Refactoring would split business logic across
# 3+ modules with tight cohesion. Deferred to v0.4.0.

import io
import json
import logging
import zipfile
import tempfile
from datetime import datetime, timezone
from pathlib import Path
from typing import IO

from app.mission.models import Mission
from app.mission.storage import load_mission_v2, load_mission_timeline
from app.mission.exporter import (
    generate_timeline_export,
    TimelineExportFormat,
)
from app.services.route_manager import RouteManager
from app.services.poi_manager import POIManager

logger = logging.getLogger(__name__)


class ExportPackageError(RuntimeError):
    """Raised when mission package export fails."""


def generate_mission_combined_csv(
    mission: Mission, output_path: str | None = None
) -> bytes | None:
    """Generate combined CSV timeline for all legs in mission.

    Combines all leg timelines into a single CSV with leg boundaries marked.
    """
    import csv

    if output_path:
        f = open(output_path, "w", newline="", encoding="utf-8")
    else:
        f = io.StringIO()

    try:
        writer = csv.writer(f)

        # Write header
        writer.writerow(
            [
                "Mission",
                mission.name,
                "Total Legs",
                len(mission.legs),
                "Generated",
                datetime.now(timezone.utc).isoformat(),
            ]
        )
        writer.writerow([])  # Blank line

        # Write combined timeline
        writer.writerow(["Leg ID", "Leg Name", "Event Time", "Event Type", "Details"])

        for leg in mission.legs:
            try:
                # Load timeline for this leg
                timeline = load_mission_timeline(leg.id)
                if not timeline:
                    continue

                # Add leg boundary marker
                writer.writerow([leg.id, leg.name, "LEG START", "---", "---"])

                # Add timeline segments
                for segment in timeline.segments:
                    start_time = (
                        segment.start_time.isoformat()
                        if isinstance(segment.start_time, datetime)
                        else segment.start_time
                    )
                    duration = (
                        (segment.end_time - segment.start_time).total_seconds()
                        if isinstance(segment.end_time, datetime)
                        else 0
                    )

                    reasons = ", ".join(segment.reasons) if segment.reasons else "---"

                    writer.writerow(
                        [
                            leg.id,
                            leg.name,
                            start_time,
                            segment.status.value,
                            f"States: X={segment.x_state.value}, Ka={segment.ka_state.value}, Ku={segment.ku_state.value} | Duration: {duration}s | Reasons: {reasons}",
                        ]
                    )

                # Add timeline advisories
                for advisory in timeline.advisories:
                    timestamp = (
                        advisory.timestamp.isoformat()
                        if isinstance(advisory.timestamp, datetime)
                        else advisory.timestamp
                    )
                    writer.writerow(
                        [
                            leg.id,
                            leg.name,
                            timestamp,
                            f"ADVISORY ({advisory.event_type})",
                            f"[{advisory.severity.upper()}] {advisory.message}",
                        ]
                    )

                writer.writerow([leg.id, leg.name, "LEG END", "---", "---"])
                writer.writerow([])  # Blank line between legs

            except Exception as e:
                logger.error(f"Failed to include leg {leg.id} in combined CSV: {e}")

    finally:
        if output_path:
            f.close()

    if not output_path:
        return f.getvalue().encode("utf-8")
    return None


def generate_mission_combined_pptx(
    mission: Mission,
    route_manager: RouteManager | None = None,
    poi_manager: POIManager | None = None,
    output_path: str | None = None,
    map_cache: dict[str, bytes] | None = None,
) -> bytes | None:
    """Generate combined PPTX slides for entire mission.

    Creates presentation with:
    - Title slide (mission overview)
    - All slides from each leg (map + timeline tables) - generated using shared pptx_builder module

    This function has been refactored to use the shared pptx_builder module,
    eliminating code duplication with exporter/__main__.py::generate_pptx_export().

    Args:
        map_cache: Optional cache for generated maps (route_id -> bytes)
    """
    import io

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        logger.error("python-pptx not installed")
        return b""

    # Import shared functions
    from pathlib import Path

    from app.mission.exporter.pptx_builder import add_mission_slides_to_presentation
    from app.mission.exporter.pptx_styling import (
        TEXT_BLACK,
        add_footer_bar,
        add_header_bar,
        add_logo,
    )

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.62)

    # Logo path
    logo_path = Path(__file__).parent.parent.joinpath("assets").joinpath("logo.png")

    # Mission metadata
    mission_id = mission.id
    mission_name = mission.name
    organization = mission.description if mission.description else "Organization"
    leg_count = len(mission.legs)

    # Title slide with styling
    blank_slide_layout = prs.slide_layouts[6]
    title_slide = prs.slides.add_slide(blank_slide_layout)

    # Add header and footer bars
    add_header_bar(title_slide, 0, 0, 10, 0.15)
    add_footer_bar(title_slide, 0, 5.47, 10, 0.15)

    # Add logo if available
    add_logo(title_slide, logo_path, 0.2, 0.02, 0.6, 0.6)

    # Add mission title
    title_box = title_slide.shapes.add_textbox(
        Inches(1.5), Inches(2.0), Inches(7.0), Inches(1.0)
    )
    text_frame = title_box.text_frame
    text_frame.text = mission_name

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(28)
    paragraph.font.bold = True
    paragraph.font.color.rgb = TEXT_BLACK

    # Add mission ID
    id_box = title_slide.shapes.add_textbox(
        Inches(1.5), Inches(3.0), Inches(7.0), Inches(0.5)
    )
    id_frame = id_box.text_frame
    id_frame.text = f"Mission ID: {mission_id}"

    id_paragraph = id_frame.paragraphs[0]
    id_paragraph.alignment = PP_ALIGN.CENTER
    id_paragraph.font.size = Pt(14)
    id_paragraph.font.color.rgb = TEXT_BLACK

    # Add leg count and organization
    info_box = title_slide.shapes.add_textbox(
        Inches(1.5), Inches(3.5), Inches(7.0), Inches(0.5)
    )
    info_frame = info_box.text_frame
    info_frame.text = f"{leg_count} Leg{'s' if leg_count != 1 else ''} | {organization}"

    info_paragraph = info_frame.paragraphs[0]
    info_paragraph.alignment = PP_ALIGN.CENTER
    info_paragraph.font.size = Pt(14)
    info_paragraph.font.color.rgb = TEXT_BLACK

    # For each leg, generate slides using shared builder
    for leg_idx, leg in enumerate(mission.legs):
        # Load timeline for this leg
        leg_timeline = load_mission_timeline(leg.id)
        if not leg_timeline:
            logger.warning(
                f"No timeline found for leg {leg.id}, adding summary slide only"
            )
        else:
            # Log timeline start time for debugging
            if leg_timeline.segments:
                first_segment_start = leg_timeline.segments[0].start_time
                logger.info(
                    f"Loaded timeline for leg {leg.id}: "
                    f"first segment starts at {first_segment_start}, "
                    f"leg.adjusted_departure_time={leg.adjusted_departure_time}"
                )

        if not leg_timeline:
            # Add a summary slide for this leg
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title_shape = slide.shapes.title
            content = slide.placeholders[1]
            title_shape.text = leg.name
            content.text = f"Leg ID: {leg.id}\nDescription: {leg.description or 'N/A'}\n\nNo timeline data available."
            continue

        try:
            # Add slides for this leg directly to main presentation
            add_mission_slides_to_presentation(
                prs=prs,
                mission=leg,
                timeline=leg_timeline,
                parent_mission_id=mission.id,
                route_manager=route_manager,
                poi_manager=poi_manager,
                logo_path=logo_path,
                map_cache=map_cache,
            )

        except Exception as e:
            logger.error(f"Failed to generate PPTX slides for leg {leg.id}: {e}")
            # Add error slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title_shape = slide.shapes.title
            content = slide.placeholders[1]
            title_shape.text = f"{leg.name} - Export Error"
            content.text = f"Leg ID: {leg.id}\n\nFailed to generate timeline: {str(e)}"

    # Save to bytes or file
    if output_path:
        prs.save(output_path)
        return None

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()


def _add_mission_metadata_to_zip(
    zf: zipfile.ZipFile, mission: Mission, manifest_files: dict
):
    """Add mission.json and leg JSON files to zip archive.

    Args:
        zf: ZipFile to add files to
        mission: Mission object to export
        manifest_files: Manifest dictionary to update
    """
    # Add mission metadata
    zf.writestr("mission.json", json.dumps(mission.model_dump(), indent=2, default=str))
    logger.info(f"Added mission.json for {mission.id}")

    # Add leg JSON files
    for leg in mission.legs:
        leg_json = json.dumps(leg.model_dump(), indent=2, default=str)
        leg_path = f"legs/{leg.id}.json"
        zf.writestr(leg_path, leg_json)
        manifest_files["legs"].append(leg_path)
        logger.info(f"Added leg file: {leg_path}")


def _add_route_kmls_to_zip(
    zf: zipfile.ZipFile,
    mission: Mission,
    route_manager: RouteManager | None,
    manifest_files: dict,
):
    """Add route KML files for all legs to zip archive.

    Args:
        zf: ZipFile to add files to
        mission: Mission object with legs
        route_manager: RouteManager instance for fetching routes
        manifest_files: Manifest dictionary to update
    """
    if not route_manager:
        return

    for leg in mission.legs:
        if leg.route_id:
            try:
                route = route_manager.get_route(leg.route_id)
                if route:
                    # Try to read the KML file from disk
                    routes_dir = Path(route_manager.routes_dir)
                    kml_file = routes_dir / f"{leg.route_id}.kml"
                    if kml_file.exists():
                        with open(kml_file, "rb") as f:
                            kml_content = f.read()
                        route_path = f"routes/{leg.route_id}.kml"
                        zf.writestr(route_path, kml_content)
                        manifest_files["routes"].append(route_path)
                        logger.info(f"Added route KML: {route_path}")
                    else:
                        logger.warning(
                            f"Route KML file not found at {kml_file} for leg {leg.id}"
                        )
                else:
                    logger.warning(f"Route {leg.route_id} not found for leg {leg.id}")
            except Exception as e:
                logger.error(f"Failed to add route KML for leg {leg.id}: {e}")


def _add_pois_to_zip(
    zf: zipfile.ZipFile,
    mission: Mission,
    poi_manager: POIManager | None,
    manifest_files: dict,
):
    """Add POI data (leg-specific and satellites) to zip archive.

    Args:
        zf: ZipFile to add files to
        mission: Mission object with legs
        poi_manager: POIManager instance for fetching POIs
        manifest_files: Manifest dictionary to update
    """
    if not poi_manager:
        return

    # Track all POI IDs we've seen to collect satellite POIs at the end
    all_poi_ids = set()

    for leg in mission.legs:
        try:
            # Get POIs associated with this leg's route and mission
            leg_pois = []

            # Get all POIs for the mission once to optimize
            all_mission_pois = poi_manager.list_pois(mission_id=mission.id)

            # Filter for POIs that are either mission-scoped (no route_id) or specific to this leg's route
            for poi in all_mission_pois:
                if poi.route_id is None or poi.route_id == leg.route_id:
                    leg_pois.append(poi)

            # Track POI IDs
            for poi in leg_pois:
                all_poi_ids.add(poi.id)

            if leg_pois:
                pois_data = {
                    "leg_id": leg.id,
                    "mission_id": mission.id,
                    "route_id": leg.route_id,
                    "pois": [poi.model_dump(mode="json") for poi in leg_pois],
                    "count": len(leg_pois),
                }
                poi_path = f"pois/{leg.id}-pois.json"
                zf.writestr(poi_path, json.dumps(pois_data, indent=2, default=str))
                manifest_files["pois"].append(poi_path)
                logger.info(f"Added POI data: {poi_path} with {len(leg_pois)} POIs")
        except Exception as e:
            logger.error(f"Failed to add POI data for leg {leg.id}: {e}")

    # Export satellite POIs (category="satellite") separately
    try:
        all_pois = poi_manager.list_pois()
        satellite_pois = [poi for poi in all_pois if poi.category == "satellite"]

        if satellite_pois:
            satellite_data = {
                "type": "satellite_pois",
                "pois": [poi.model_dump(mode="json") for poi in satellite_pois],
                "count": len(satellite_pois),
            }
            satellite_path = "pois/satellites.json"
            zf.writestr(
                satellite_path,
                json.dumps(satellite_data, indent=2, default=str),
            )
            manifest_files["pois"].append(satellite_path)
            logger.info(
                f"Added satellite POI data: {satellite_path} with {len(satellite_pois)} satellites"
            )
    except Exception as e:
        logger.error(f"Failed to add satellite POI data: {e}")


def _add_per_leg_exports_to_zip(
    zf: zipfile.ZipFile,
    mission: Mission,
    route_manager: RouteManager | None,
    poi_manager: POIManager | None,
    manifest_files: dict,
    map_cache: dict[str, bytes] | None = None,
):
    """Generate and add per-leg exports (CSV, PPTX) to zip archive.

    Args:
        zf: ZipFile to add files to
        mission: Mission object with legs
        route_manager: RouteManager instance
        poi_manager: POIManager instance
        manifest_files: Manifest dictionary to update
        map_cache: Optional cache for generated maps (route_id -> bytes)
    """
    for leg in mission.legs:
        # Load timeline for this specific leg
        leg_timeline = load_mission_timeline(leg.id)
        if not leg_timeline:
            logger.warning(
                f"No timeline found for leg {leg.id}, skipping exports for this leg"
            )
            continue

        # Log timeline start time for debugging adjusted departure times
        if leg_timeline.segments:
            first_segment_start = leg_timeline.segments[0].start_time
            logger.info(
                f"Exporting leg {leg.id}: "
                f"timeline starts at {first_segment_start}, "
                f"leg.adjusted_departure_time={leg.adjusted_departure_time}"
            )

        try:
            # CSV export
            csv_export = generate_timeline_export(
                export_format=TimelineExportFormat.CSV,
                mission=leg,  # Pass the leg as mission (MissionLeg has same fields)
                timeline=leg_timeline,
            )
            csv_path = f"exports/legs/{leg.id}/timeline.csv"
            zf.writestr(csv_path, csv_export.content)
            manifest_files["per_leg_exports"].append(csv_path)
            logger.info(f"Generated CSV for leg {leg.id}")
        except Exception as e:
            logger.error(f"Failed to generate CSV for leg {leg.id}: {e}")

        try:
            # PPTX export
            pptx_export = generate_timeline_export(
                export_format=TimelineExportFormat.PPTX,
                mission=leg,
                timeline=leg_timeline,
                parent_mission_id=mission.id,
                route_manager=route_manager,
                poi_manager=poi_manager,
                map_cache=map_cache,
            )
            pptx_path = f"exports/legs/{leg.id}/slides.pptx"
            zf.writestr(pptx_path, pptx_export.content)
            manifest_files["per_leg_exports"].append(pptx_path)
            logger.info(f"Generated PPTX for leg {leg.id}")
        except Exception as e:
            logger.error(f"Failed to generate PPTX for leg {leg.id}: {e}")


def _add_combined_mission_exports_to_zip(
    zf: zipfile.ZipFile,
    mission: Mission,
    route_manager: RouteManager | None,
    poi_manager: POIManager | None,
    manifest_files: dict,
    map_cache: dict[str, bytes] | None = None,
):
    """Generate and add combined mission-level exports (CSV, PPTX) to zip archive.

    Args:
        zf: ZipFile to add files to
        mission: Mission object
        route_manager: RouteManager instance
        poi_manager: POIManager instance
        manifest_files: Manifest dictionary to update
        map_cache: Optional cache for generated maps (route_id -> bytes)
    """
    try:
        logger.info(f"Generating combined mission-level exports for {mission.id}")

        # Combined CSV - stream to temp file
        with tempfile.NamedTemporaryFile(delete=True) as tmp_csv:
            generate_mission_combined_csv(mission, output_path=tmp_csv.name)
            zf.write(tmp_csv.name, "exports/mission/mission-timeline.csv")
            manifest_files["mission_exports"].append(
                "exports/mission/mission-timeline.csv"
            )

        # Combined PPTX - stream to temp file
        with tempfile.NamedTemporaryFile(delete=True, suffix=".pptx") as tmp_pptx:
            generate_mission_combined_pptx(
                mission,
                route_manager=route_manager,
                poi_manager=poi_manager,
                output_path=tmp_pptx.name,
                map_cache=map_cache,
            )
            zf.write(tmp_pptx.name, "exports/mission/mission-slides.pptx")
            manifest_files["mission_exports"].append(
                "exports/mission/mission-slides.pptx"
            )

        logger.info("Combined mission-level exports complete")

    except Exception as e:
        logger.error(f"Failed to generate combined mission exports: {e}")


def _create_export_manifest(mission: Mission, manifest_files: dict) -> dict:
    """Create the manifest.json content for the export package.

    Args:
        mission: Mission object
        manifest_files: Dictionary with file lists by category

    Returns:
        Manifest dictionary
    """
    return {
        "version": "2.0",
        "mission_id": mission.id,
        "mission_name": mission.name,
        "mission_description": mission.description or "",
        "leg_count": len(mission.legs),
        "exported_at": datetime.now(timezone.utc).isoformat(),
        "file_structure": {
            "mission_data": manifest_files["mission_data"],
            "legs": manifest_files["legs"],
            "routes": manifest_files["routes"],
            "pois": manifest_files["pois"],
            "mission_exports": manifest_files["mission_exports"],
            "per_leg_exports": manifest_files["per_leg_exports"],
        },
        "statistics": {
            "total_files": (
                len(manifest_files["mission_data"])
                + len(manifest_files["legs"])
                + len(manifest_files["routes"])
                + len(manifest_files["pois"])
                + len(manifest_files["mission_exports"])
                + len(manifest_files["per_leg_exports"])
            ),
            "leg_files": len(manifest_files["legs"]),
            "route_files": len(manifest_files["routes"]),
            "poi_files": len(manifest_files["pois"]),
            "mission_export_files": len(manifest_files["mission_exports"]),
            "per_leg_export_files": len(manifest_files["per_leg_exports"]),
        },
    }


def export_mission_package(
    mission_id: str,
    route_manager: RouteManager,
    poi_manager: POIManager,
) -> IO[bytes]:
    """Export complete mission as zip archive.

    Package structure:
        mission-{id}.zip
        ├── mission.json
        ├── manifest.json
        ├── legs/
        │   ├── {leg-id-1}.json
        │   └── {leg-id-2}.json
        ├── routes/
        │   ├── {route-id-1}.kml
        │   └── {route-id-2}.kml
        ├── pois/
        │   ├── {leg-id-1}-pois.json
        │   └── {leg-id-2}-pois.json
        ├── exports/
        │   ├── mission/
        │   │   ├── mission-timeline.csv
        │   │   ├── mission-timeline.xlsx
        │   │   ├── mission-slides.pptx
        │   │   └── mission-report.pdf
        │   └── legs/
        │       ├── {leg-id-1}/
        │       │   ├── timeline.csv
        │       │   ├── timeline.xlsx
        │       │   ├── slides.pptx
        │       │   └── report.pdf
        │       └── {leg-id-2}/
        │           ├── timeline.csv
        │           ├── timeline.xlsx
        │           ├── slides.pptx
        │           └── report.pdf

    Args:
        mission_id: Mission to export
        route_manager: RouteManager instance for fetching route KML files
        poi_manager: POIManager instance for fetching mission POIs

    Returns:
        File-like object containing the zip archive. Caller must close it to delete the temp file.
    """
    mission = load_mission_v2(mission_id)

    if not mission:
        raise ExportPackageError(f"Mission {mission_id} not found")

    # Create a temporary file for the zip archive
    # delete=True ensures it's deleted when closed by the caller (StreamingResponse)
    zip_temp = tempfile.NamedTemporaryFile(delete=True)

    manifest_files = {
        "mission_data": ["mission.json", "manifest.json"],
        "legs": [],
        "routes": [],
        "pois": [],
        "mission_exports": [],
        "per_leg_exports": [],
    }

    # Create map cache for this export operation to avoid regenerating same maps
    map_cache: dict[str, bytes] = {}

    try:
        with zipfile.ZipFile(zip_temp, "w", zipfile.ZIP_DEFLATED) as zf:
            # Add mission metadata and leg files
            _add_mission_metadata_to_zip(zf, mission, manifest_files)

            # Add route KML files
            _add_route_kmls_to_zip(zf, mission, route_manager, manifest_files)

            # Add POI data (leg-specific and satellites)
            _add_pois_to_zip(zf, mission, poi_manager, manifest_files)

            # Generate and add per-leg exports (will populate map_cache)
            _add_per_leg_exports_to_zip(
                zf, mission, route_manager, poi_manager, manifest_files, map_cache
            )

            # Generate and add combined mission-level exports (will reuse cached maps)
            _add_combined_mission_exports_to_zip(
                zf, mission, route_manager, poi_manager, manifest_files, map_cache
            )

            # Create and add manifest
            manifest = _create_export_manifest(mission, manifest_files)
            manifest_json = json.dumps(manifest, indent=2)
            zf.writestr("manifest.json", manifest_json)
            logger.info(
                f"Created manifest.json with {manifest['statistics']['total_files']} files"
            )

    except Exception:
        # If anything fails, close and delete the temp file
        zip_temp.close()
        raise

    # Seek to beginning so it can be read
    zip_temp.seek(0)
    return zip_temp
