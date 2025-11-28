"""Mission package export utilities for creating portable mission archives."""

import copy
import hashlib
import io
import json
import logging
import zipfile
import tempfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional, Generator, IO, Any

from app.mission.models import Mission, MissionLeg
from app.mission.storage import load_mission_v2, get_mission_path, load_mission_timeline
from app.mission.exporter import generate_timeline_export, TimelineExportFormat, ExportGenerationError
from app.services.route_manager import RouteManager
from app.services.poi_manager import POIManager

logger = logging.getLogger(__name__)

# Constants
EXCEL_SHEET_NAME_MAX_LENGTH = 31  # Excel's maximum sheet name length




class ExportPackageError(RuntimeError):
    """Raised when mission package export fails."""


def _create_mission_summary_sheet(wb, mission: Mission):
    """Create mission summary sheet with metadata and leg index.
    
    Args:
        wb: Workbook to add sheet to
        mission: Mission object with metadata
        
    Returns:
        Worksheet with mission summary
    """
    ws = wb.active
    ws.title = "Mission Summary"
    ws.append(["Mission Name", mission.name])
    ws.append(["Mission ID", mission.id])
    ws.append(["Total Legs", len(mission.legs)])
    ws.append(["Generated", datetime.now(timezone.utc).isoformat()])
    ws.append([])
    ws.append(["Leg #", "Leg ID", "Leg Name", "Description"])

    for idx, leg in enumerate(mission.legs, 1):
        ws.append([idx, leg.id, leg.name, leg.description or ""])

    # Set column widths for better readability
    ws.column_dimensions["A"].width = 10
    ws.column_dimensions["B"].width = 25
    ws.column_dimensions["C"].width = 30
    ws.column_dimensions["D"].width = 50
    
    return ws


def _copy_worksheet_content(source_sheet, target_sheet):
    """Copy all content from source worksheet to target worksheet.
    
    Args:
        source_sheet: Source worksheet to copy from
        target_sheet: Target worksheet to copy to
    """
    # Copy cell values, styles, and merged cells
    for row in source_sheet.iter_rows():
        for cell in row:
            new_cell = target_sheet[cell.coordinate]
            new_cell.value = cell.value

            # Copy cell formatting
            if cell.has_style:
                new_cell.font = copy.copy(cell.font)
                new_cell.border = copy.copy(cell.border)
                new_cell.fill = copy.copy(cell.fill)
                new_cell.number_format = copy.copy(cell.number_format)
                new_cell.protection = copy.copy(cell.protection)
                new_cell.alignment = copy.copy(cell.alignment)

    # Copy column widths
    for col_letter, col_dim in source_sheet.column_dimensions.items():
        target_sheet.column_dimensions[col_letter].width = col_dim.width

    # Copy row heights
    for row_num, row_dim in source_sheet.row_dimensions.items():
        target_sheet.row_dimensions[row_num].height = row_dim.height

    # Copy merged cells
    for merged_cell_range in source_sheet.merged_cells.ranges:
        target_sheet.merge_cells(str(merged_cell_range))

    # Copy images if any
    for image in source_sheet._images:
        new_image = copy.copy(image)
        target_sheet.add_image(new_image, image.anchor)


def _add_error_sheet(wb, leg_idx: int, leg, error_message: str):
    """Add an error sheet for a failed leg export.
    
    Args:
        wb: Workbook to add sheet to
        leg_idx: Leg index (0-based)
        leg: Leg object
        error_message: Error message to display
        
    Returns:
        Worksheet with error information
    """
    ws_error = wb.create_sheet(title=f"L{leg_idx + 1} Error")
    ws_error.append(["Leg Name", leg.name])
    ws_error.append(["Leg ID", leg.id])
    ws_error.append([])
    ws_error.append(["Error", error_message])
    return ws_error


def _process_leg_xlsx_export(
    wb,
    leg,
    leg_idx: int,
    mission_id: str,
    route_manager: RouteManager | None,
    poi_manager: POIManager | None,
):
    """Process a single leg's XLSX export and add sheets to workbook.
    
    Args:
        wb: Workbook to add sheets to
        leg: Leg object to process
        leg_idx: Leg index (0-based)
        mission_id: Parent mission ID
        route_manager: RouteManager instance
        poi_manager: POIManager instance
    """
    from openpyxl import load_workbook
    
    # Load timeline for this leg
    leg_timeline = load_mission_timeline(leg.id)
    if not leg_timeline:
        logger.warning(
            f"No timeline found for leg {leg.id}, adding summary sheet only"
        )
        # Add a simple sheet for this leg
        ws_leg = wb.create_sheet(title=f"Leg {leg_idx + 1} - {leg.name[:20]}")
        ws_leg.append(["Leg Name", leg.name])
        ws_leg.append(["Leg ID", leg.id])
        ws_leg.append(["Description", leg.description or "N/A"])
        ws_leg.append([])
        ws_leg.append(["Status", "No timeline data available"])
        return

    try:
        # Generate full XLSX export for this leg
        leg_xlsx_bytes = generate_timeline_export(
            export_format=TimelineExportFormat.XLSX,
            mission=leg,
            timeline=leg_timeline,
            parent_mission_id=mission_id,
            route_manager=route_manager,
            poi_manager=poi_manager,
        )

        # Load the leg workbook
        leg_wb = load_workbook(io.BytesIO(leg_xlsx_bytes.content))

        # Copy each sheet from leg workbook to main workbook with leg prefix
        for sheet_name in leg_wb.sheetnames:
            leg_sheet = leg_wb[sheet_name]

            # Create new sheet name with leg prefix
            # Truncate leg name to fit Excel's sheet name limit
            leg_prefix = f"L{leg_idx + 1}"
            if sheet_name == "Summary":
                new_sheet_name = f"{leg_prefix} {leg.name[:20]}"
            else:
                new_sheet_name = f"{leg_prefix} {sheet_name[:25]}"

            # Ensure sheet name is unique and under limit
            new_sheet_name = new_sheet_name[:EXCEL_SHEET_NAME_MAX_LENGTH]

            # Create new sheet
            new_sheet = wb.create_sheet(title=new_sheet_name)

            # Copy all content from leg sheet to new sheet
            _copy_worksheet_content(leg_sheet, new_sheet)

    except ExportGenerationError as e:
        logger.error(f"Failed to generate XLSX export for leg {leg.id}: {e}")
        _add_error_sheet(wb, leg_idx, leg, str(e))

    except (ValueError, KeyError) as e:
        logger.error(f"Sheet manipulation error for leg {leg.id}: {e}", exc_info=True)
        _add_error_sheet(wb, leg_idx, leg, f"Processing Error: {str(e)}")

    except zipfile.BadZipFile as e:
        logger.error(f"Invalid XLSX data generated for leg {leg.id}: {e}", exc_info=True)
        _add_error_sheet(wb, leg_idx, leg, "Generated Excel file was invalid")

    except Exception as e:
        logger.error(f"Unexpected error combining XLSX for leg {leg.id}: {e}", exc_info=True)
        _add_error_sheet(wb, leg_idx, leg, f"Unexpected Error: {str(e)}")


def generate_mission_combined_csv(mission: Mission, output_path: str | None = None) -> bytes | None:
    """Generate combined CSV timeline for all legs in mission.

    Combines all leg timelines into a single CSV with leg boundaries marked.
    """
    import csv

    if output_path:
        f = open(output_path, "w", newline="", encoding="utf-8")
    else:
        f = io.StringIO()

    try:
        writer = csv.writer(f)

        # Write header
        writer.writerow(
            [
                "Mission",
                mission.name,
                "Total Legs",
                len(mission.legs),
                "Generated",
                datetime.now(timezone.utc).isoformat(),
            ]
        )
        writer.writerow([])  # Blank line

        # Write combined timeline
        writer.writerow(["Leg ID", "Leg Name", "Event Time", "Event Type", "Details"])

        for leg in mission.legs:
            try:
                # Load timeline for this leg
                timeline = load_mission_timeline(leg.id)
                if not timeline:
                    continue

                # Add leg boundary marker
                writer.writerow([leg.id, leg.name, "LEG START", "---", "---"])

                # Add timeline segments
                for segment in timeline.segments:
                    start_time = (
                        segment.start_time.isoformat()
                        if isinstance(segment.start_time, datetime)
                        else segment.start_time
                    )
                    end_time = (
                        segment.end_time.isoformat()
                        if isinstance(segment.end_time, datetime)
                        else segment.end_time
                    )
                    duration = (
                        (segment.end_time - segment.start_time).total_seconds()
                        if isinstance(segment.end_time, datetime)
                        else 0
                    )

                    reasons = ", ".join(segment.reasons) if segment.reasons else "---"

                    writer.writerow(
                        [
                            leg.id,
                            leg.name,
                            start_time,
                            segment.status.value,
                            f"States: X={segment.x_state.value}, Ka={segment.ka_state.value}, Ku={segment.ku_state.value} | Duration: {duration}s | Reasons: {reasons}",
                        ]
                    )

                # Add timeline advisories
                for advisory in timeline.advisories:
                    timestamp = (
                        advisory.timestamp.isoformat()
                        if isinstance(advisory.timestamp, datetime)
                        else advisory.timestamp
                    )
                    writer.writerow(
                        [
                            leg.id,
                            leg.name,
                            timestamp,
                            f"ADVISORY ({advisory.event_type})",
                            f"[{advisory.severity.upper()}] {advisory.message}",
                        ]
                    )

                writer.writerow([leg.id, leg.name, "LEG END", "---", "---"])
                writer.writerow([])  # Blank line between legs

            except Exception as e:
                logger.error(f"Failed to include leg {leg.id} in combined CSV: {e}")

    finally:
        if output_path:
            f.close()

    if not output_path:
        return f.getvalue().encode("utf-8")
    return None


def generate_mission_combined_xlsx(
    mission: Mission,
    route_manager: RouteManager | None = None,
    poi_manager: POIManager | None = None,
    output_path: str | None = None,
) -> bytes | None:
    """Generate combined XLSX timeline with one sheet per leg plus summary.

    Creates workbook with:
    - Mission Summary sheet (overview + leg index)
    - Full sheets from each leg (Summary, Timeline, Advisories, Statistics)
    """
    try:
        from openpyxl import Workbook
    except ImportError:
        logger.error("openpyxl not installed")
        return b""

    try:
        # Create main workbook
        wb = Workbook()

        # Create mission summary sheet
        _create_mission_summary_sheet(wb, mission)

        # Process each leg and import its sheets
        for leg_idx, leg in enumerate(mission.legs):
            _process_leg_xlsx_export(
                wb, leg, leg_idx, mission.id, route_manager, poi_manager
            )

        # Save to bytes or file
        if output_path:
            wb.save(output_path)
            return None
        
        output = io.BytesIO()
        wb.save(output)
        output.seek(0)
        return output.read()

    except Exception as e:
        logger.error(f"Failed to generate combined XLSX: {e}", exc_info=True)
        # Return a basic error workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Error"
        ws.append(["Error generating combined Excel file"])
        ws.append([str(e)])
        
        if output_path:
            wb.save(output_path)
            return None
            
        output = io.BytesIO()
        wb.save(output)
        output.seek(0)
        return output.read()


def generate_mission_combined_pptx(
    mission: Mission,
    route_manager: RouteManager | None = None,
    poi_manager: POIManager | None = None,
    output_path: str | None = None,
) -> bytes | None:
    """Generate combined PPTX slides for entire mission.

    Creates presentation with:
    - Title slide (mission overview)
    - All slides from each leg (map + timeline tables) - generated using the same code as individual leg exports
    - Summary slide
    """
    import io

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        logger.error("python-pptx not installed")
        return b""

    # Import the helper function from exporter that generates slides for a leg
    from app.mission.exporter import (
        _generate_route_map,
        _segment_rows,
        _generate_timeline_chart,
    )
    from app.mission.exporter import TRANSPORT_DISPLAY, STATE_COLUMNS, Transport
    from openpyxl.styles import PatternFill
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = mission.name
    subtitle.text = f"{len(mission.legs)} Legs"

    # For each leg, generate slides directly (same logic as generate_pptx_export)
    for leg_idx, leg in enumerate(mission.legs):
        # Load timeline for this leg
        leg_timeline = load_mission_timeline(leg.id)
        if not leg_timeline:
            logger.warning(
                f"No timeline found for leg {leg.id}, adding summary slide only"
            )
            # Add a summary slide for this leg
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title_shape = slide.shapes.title
            content = slide.placeholders[1]
            title_shape.text = leg.name
            content.text = f"Leg ID: {leg.id}\nDescription: {leg.description or 'N/A'}\n\nNo timeline data available."
            continue

        try:
            # Generate slides for this leg using the same logic as generate_pptx_export
            # Slide 1: Route Map
            blank_slide_layout = prs.slide_layouts[6]
            slide1 = prs.slides.add_slide(blank_slide_layout)

            # Generate map image
            try:
                map_image_bytes = _generate_route_map(
                    leg_timeline,
                    leg,
                    parent_mission_id=mission.id,
                    route_manager=route_manager,
                    poi_manager=poi_manager,
                )
                map_image_stream = io.BytesIO(map_image_bytes)

                # Add image to slide
                left = Inches(0.5)
                top = Inches(0.5)
                width = Inches(9)
                height = Inches(6.5)

                slide1.shapes.add_picture(
                    map_image_stream, left, top, width=width, height=height
                )

                # Add title
                title_box = slide1.shapes.add_textbox(
                    Inches(0.5), Inches(0.1), Inches(9), Inches(0.5)
                )
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = f"{leg.name} - Route Map"
                p.font.bold = True
                p.font.size = Pt(24)
                p.alignment = PP_ALIGN.CENTER

            except Exception as e:
                logger.error(
                    f"Failed to add map to PPTX for leg {leg.id}: {e}", exc_info=True
                )
                textbox = slide1.shapes.add_textbox(
                    Inches(1), Inches(1), Inches(8), Inches(1)
                )
                tf = textbox.text_frame
                tf.text = f"Map generation failed: {str(e)}"

            # Slide 2+: Timeline Table (Paginated) - using same logic as generate_pptx_export
            timeline_df = _segment_rows(leg_timeline, leg)

            if not timeline_df.empty:
                from pptx.util import Pt
                from pptx.dml.color import RGBColor

                columns_to_show = [
                    "Segment #",
                    "Status",
                    "Start Time",
                    "End Time",
                    "Duration",
                    TRANSPORT_DISPLAY[Transport.X],
                    TRANSPORT_DISPLAY[Transport.KA],
                    TRANSPORT_DISPLAY[Transport.KU],
                    "Reasons",
                ]

                # Pagination settings
                ROWS_PER_SLIDE = 10
                MIN_ROWS_LAST_SLIDE = 3

                records = timeline_df.to_dict("records")
                total_rows = len(records)

                chunks = []
                if total_rows <= ROWS_PER_SLIDE:
                    chunks.append(timeline_df)
                else:
                    current_idx = 0
                    while current_idx < total_rows:
                        end_idx = min(current_idx + ROWS_PER_SLIDE, total_rows)
                        chunks.append(timeline_df.iloc[current_idx:end_idx])
                        current_idx = end_idx

                    if len(chunks) > 1:
                        last_chunk = chunks[-1]
                        if len(last_chunk) < MIN_ROWS_LAST_SLIDE:
                            needed = MIN_ROWS_LAST_SLIDE - len(last_chunk)
                            base_chunks = chunks[:-2]
                            remainder_count = len(chunks[-2]) + len(chunks[-1])
                            split_idx = total_rows - remainder_count
                            second_last_len = remainder_count - MIN_ROWS_LAST_SLIDE

                            chunk_second_last = timeline_df.iloc[
                                split_idx : split_idx + second_last_len
                            ]
                            chunk_last = timeline_df.iloc[split_idx + second_last_len :]

                            chunks = base_chunks + [chunk_second_last, chunk_last]

                for chunk_idx, chunk in enumerate(chunks):
                    slide = prs.slides.add_slide(blank_slide_layout)

                    # Add title
                    title_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.1), Inches(9), Inches(0.5)
                    )
                    tf = title_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = (
                        f"{leg.name} - Timeline"
                        if chunk_idx == 0
                        else f"{leg.name} - Timeline (cont.)"
                    )
                    p.font.bold = True
                    p.font.size = Pt(24)
                    p.alignment = PP_ALIGN.CENTER

                    rows = len(chunk) + 1
                    cols = len(columns_to_show)

                    left = Inches(0.5)
                    top = Inches(1.0)
                    width = Inches(9.0)
                    height = Inches(1.0)

                    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
                    table = shape.table

                    col_weights = [0.6, 1.0, 1.75, 1.75, 1.0, 1.0, 1.0, 1.0, 1.5]
                    total_weight = sum(col_weights)
                    for i, weight in enumerate(col_weights):
                        table.columns[i].width = int(width * (weight / total_weight))

                    # Header row
                    for i, col_name in enumerate(columns_to_show):
                        cell = table.cell(0, i)
                        cell.text = col_name
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(51, 102, 178)
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
                        paragraph.font.bold = True
                        paragraph.font.size = Pt(8)
                        paragraph.alignment = PP_ALIGN.CENTER

                    # Data rows
                    for row_idx, (_, row_data) in enumerate(chunk.iterrows(), start=1):
                        bad_transports = []
                        for col_name in STATE_COLUMNS:
                            val = str(
                                row_data[col_name] if row_data[col_name] else ""
                            ).lower()
                            if val in ("degraded", "warning", "offline"):
                                bad_transports.append(col_name)

                        is_critical_row = len(bad_transports) >= 2

                        for col_idx, col_name in enumerate(columns_to_show):
                            cell = table.cell(row_idx, col_idx)
                            val = str(
                                row_data[col_name]
                                if row_data[col_name] is not None
                                else ""
                            )

                            if col_name == "Status" and is_critical_row:
                                val = "CRITICAL"

                            cell.text = val

                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.font.size = Pt(7)
                                paragraph.alignment = PP_ALIGN.LEFT

                            cell.fill.solid()
                            if row_idx % 2 == 0:
                                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                            else:
                                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

                            if col_name in STATE_COLUMNS or col_name == "Status":
                                val_lower = val.lower()

                                is_sof = False
                                reasons = str(row_data.get("Reasons", "")).lower()
                                if "safety-of-flight" in reasons or "aar" in reasons:
                                    is_sof = True

                                if (
                                    col_name == "Status"
                                    and is_sof
                                    and val_lower in ("available", "nominal", "warning")
                                ):
                                    cell.text = "SOF"
                                    for paragraph in cell.text_frame.paragraphs:
                                        paragraph.font.size = Pt(7)
                                        paragraph.alignment = PP_ALIGN.LEFT

                                if (
                                    val_lower in ("degraded", "warning", "offline")
                                    or is_sof
                                ):
                                    cell.fill.solid()
                                    if is_critical_row and not is_sof:
                                        cell.fill.fore_color.rgb = RGBColor(
                                            255, 204, 204
                                        )
                                    elif val_lower in ("degraded", "warning") or is_sof:
                                        cell.fill.fore_color.rgb = RGBColor(
                                            255, 255, 204
                                        )
                                    else:
                                        cell.fill.fore_color.rgb = RGBColor(
                                            255, 204, 204
                                        )

        except Exception as e:
            logger.error(f"Failed to generate PPTX slides for leg {leg.id}: {e}")
            # Add error slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title_shape = slide.shapes.title
            content = slide.placeholders[1]
            title_shape.text = f"{leg.name} - Export Error"
            content.text = f"Leg ID: {leg.id}\n\nFailed to generate timeline: {str(e)}"

    # Save to bytes or file
    if output_path:
        prs.save(output_path)
        return None

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()


def generate_mission_combined_pdf(
    mission: Mission,
    route_manager: RouteManager | None = None,
    poi_manager: POIManager | None = None,
    output_path: str | None = None,
) -> bytes | None:
    """Generate combined PDF report for entire mission.

    Creates PDF with:
    - Cover page (mission overview)
    - Full PDF report for each leg (with maps and timelines)
    - Summary page
    """
    import io

    try:
        from PyPDF2 import PdfMerger, PdfReader
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
    except ImportError:
        logger.error("PyPDF2 or reportlab not installed")
        return b""

    merger = PdfMerger()

    try:
        # Create cover page
        cover_buffer = io.BytesIO()
        c = canvas.Canvas(cover_buffer, pagesize=letter)
        width, height = letter

        c.setFont("Helvetica-Bold", 24)
        c.drawString(100, height - 100, mission.name)
        c.setFont("Helvetica", 12)
        c.drawString(100, height - 140, f"Mission ID: {mission.id}")
        c.drawString(100, height - 160, f"Total Legs: {len(mission.legs)}")
        c.drawString(
            100,
            height - 200,
            f"Generated: {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M:%S UTC')}",
        )

        # Add table of contents
        c.setFont("Helvetica-Bold", 14)
        c.drawString(100, height - 250, "Table of Contents")
        c.setFont("Helvetica", 10)
        y_position = height - 280
        for idx, leg in enumerate(mission.legs, 1):
            c.drawString(120, y_position, f"{idx}. {leg.name}")
            y_position -= 20

        c.showPage()
        c.save()
        cover_buffer.seek(0)

        # Add cover page to merger
        merger.append(cover_buffer)

        # Generate and add each leg's full PDF
        for leg_idx, leg in enumerate(mission.legs):
            # Load timeline for this leg
            leg_timeline = load_mission_timeline(leg.id)
            if not leg_timeline:
                logger.warning(
                    f"No timeline found for leg {leg.id}, adding summary page only"
                )
                # Create a simple page for this leg
                leg_buffer = io.BytesIO()
                c = canvas.Canvas(leg_buffer, pagesize=letter)
                c.setFont("Helvetica-Bold", 18)
                c.drawString(100, height - 100, f"Leg: {leg.name}")
                c.setFont("Helvetica", 12)
                c.drawString(100, height - 140, f"Leg ID: {leg.id}")
                c.drawString(
                    100, height - 160, f"Description: {leg.description or 'N/A'}"
                )
                c.drawString(100, height - 200, "No timeline data available.")
                c.showPage()
                c.save()
                leg_buffer.seek(0)
                merger.append(leg_buffer)
                continue

            try:
                # Generate full PDF export for this leg
                leg_pdf_bytes = generate_timeline_export(
                    export_format=TimelineExportFormat.PDF,
                    mission=leg,
                    timeline=leg_timeline,
                    parent_mission_id=mission.id,
                    route_manager=route_manager,
                    poi_manager=poi_manager,
                )

                # Add section divider page before leg PDF (except for first leg)
                if leg_idx > 0:
                    divider_buffer = io.BytesIO()
                    c = canvas.Canvas(divider_buffer, pagesize=letter)
                    c.setFont("Helvetica-Bold", 20)
                    c.drawString(100, height - 100, f"Leg {leg_idx + 1}")
                    c.setFont("Helvetica-Bold", 16)
                    c.drawString(100, height - 140, leg.name)
                    c.showPage()
                    c.save()
                    divider_buffer.seek(0)
                    merger.append(divider_buffer)

                # Append the leg PDF
                leg_pdf_buffer = io.BytesIO(leg_pdf_bytes.content)
                merger.append(leg_pdf_buffer)

            except Exception as e:
                logger.error(f"Failed to combine PDF for leg {leg.id}: {e}")
                # Add error page
                error_buffer = io.BytesIO()
                c = canvas.Canvas(error_buffer, pagesize=letter)
                c.setFont("Helvetica-Bold", 18)
                c.drawString(100, height - 100, f"Leg: {leg.name} - Export Error")
                c.setFont("Helvetica", 12)
                c.drawString(100, height - 140, f"Leg ID: {leg.id}")
                c.drawString(100, height - 180, f"Error: {str(e)}")
                c.showPage()
                c.save()
                error_buffer.seek(0)
                merger.append(error_buffer)

        # Create summary page
        summary_buffer = io.BytesIO()
        c = canvas.Canvas(summary_buffer, pagesize=letter)

        c.setFont("Helvetica-Bold", 20)
        c.drawString(100, height - 100, "Mission Summary")
        c.setFont("Helvetica", 12)
        c.drawString(100, height - 140, f"Total Legs: {len(mission.legs)}")

        y_position = height - 180
        for idx, leg in enumerate(mission.legs, 1):
            c.drawString(100, y_position, f"{idx}. {leg.name} ({leg.id})")
            y_position -= 20

        c.showPage()
        c.save()
        summary_buffer.seek(0)

        # Add summary page
        merger.append(summary_buffer)

        # Write final combined PDF
        if output_path:
            merger.write(output_path)
            merger.close()
            return None

        output = io.BytesIO()
        merger.write(output)
        merger.close()
        output.seek(0)
        return output.read()

    except Exception as e:
        logger.error(f"Failed to generate combined PDF: {e}")
        # Return a basic error PDF
        output_stream = io.BytesIO()
        c = canvas.Canvas(output_stream, pagesize=letter)
        width, height = letter
        c.setFont("Helvetica-Bold", 18)
        c.drawString(100, height - 100, "PDF Generation Error")
        c.setFont("Helvetica", 12)
        c.drawString(100, height - 140, str(e))
        c.showPage()
        c.save()
        output_stream.seek(0)
        
        if output_path:
            with open(output_path, "wb") as f:
                f.write(output_stream.read())
            return None
            
        return output_stream.read()


def _add_mission_metadata_to_zip(zf: zipfile.ZipFile, mission: Mission, manifest_files: dict):
    """Add mission.json and leg JSON files to zip archive.
    
    Args:
        zf: ZipFile to add files to
        mission: Mission object to export
        manifest_files: Manifest dictionary to update
    """
    # Add mission metadata
    zf.writestr(
        "mission.json", json.dumps(mission.model_dump(), indent=2, default=str)
    )
    logger.info(f"Added mission.json for {mission.id}")

    # Add leg JSON files
    for leg in mission.legs:
        leg_json = json.dumps(leg.model_dump(), indent=2, default=str)
        leg_path = f"legs/{leg.id}.json"
        zf.writestr(leg_path, leg_json)
        manifest_files["legs"].append(leg_path)
        logger.info(f"Added leg file: {leg_path}")


def _add_route_kmls_to_zip(
    zf: zipfile.ZipFile,
    mission: Mission,
    route_manager: RouteManager | None,
    manifest_files: dict,
):
    """Add route KML files for all legs to zip archive.
    
    Args:
        zf: ZipFile to add files to
        mission: Mission object with legs
        route_manager: RouteManager instance for fetching routes
        manifest_files: Manifest dictionary to update
    """
    if not route_manager:
        return

    for leg in mission.legs:
        if leg.route_id:
            try:
                route = route_manager.get_route(leg.route_id)
                if route:
                    # Try to read the KML file from disk
                    routes_dir = Path(route_manager.routes_dir)
                    kml_file = routes_dir / f"{leg.route_id}.kml"
                    if kml_file.exists():
                        with open(kml_file, "rb") as f:
                            kml_content = f.read()
                        route_path = f"routes/{leg.route_id}.kml"
                        zf.writestr(route_path, kml_content)
                        manifest_files["routes"].append(route_path)
                        logger.info(f"Added route KML: {route_path}")
                    else:
                        logger.warning(
                            f"Route KML file not found at {kml_file} for leg {leg.id}"
                        )
                else:
                    logger.warning(
                        f"Route {leg.route_id} not found for leg {leg.id}"
                    )
            except Exception as e:
                logger.error(f"Failed to add route KML for leg {leg.id}: {e}")


def _add_pois_to_zip(
    zf: zipfile.ZipFile,
    mission: Mission,
    poi_manager: POIManager | None,
    manifest_files: dict,
):
    """Add POI data (leg-specific and satellites) to zip archive.
    
    Args:
        zf: ZipFile to add files to
        mission: Mission object with legs
        poi_manager: POIManager instance for fetching POIs
        manifest_files: Manifest dictionary to update
    """
    if not poi_manager:
        return

    # Track all POI IDs we've seen to collect satellite POIs at the end
    all_poi_ids = set()

    for leg in mission.legs:
        try:
            # Get POIs associated with this leg's route and mission
            leg_pois = []

            # Get all POIs for the mission once to optimize
            all_mission_pois = poi_manager.list_pois(mission_id=mission.id)

            # Filter for POIs that are either mission-scoped (no route_id) or specific to this leg's route
            for poi in all_mission_pois:
                if poi.route_id is None or poi.route_id == leg.route_id:
                    leg_pois.append(poi)

            # Track POI IDs
            for poi in leg_pois:
                all_poi_ids.add(poi.id)

            if leg_pois:
                pois_data = {
                    "leg_id": leg.id,
                    "mission_id": mission.id,
                    "route_id": leg.route_id,
                    "pois": [poi.model_dump(mode="json") for poi in leg_pois],
                    "count": len(leg_pois),
                }
                poi_path = f"pois/{leg.id}-pois.json"
                zf.writestr(
                    poi_path, json.dumps(pois_data, indent=2, default=str)
                )
                manifest_files["pois"].append(poi_path)
                logger.info(
                    f"Added POI data: {poi_path} with {len(leg_pois)} POIs"
                )
        except Exception as e:
            logger.error(f"Failed to add POI data for leg {leg.id}: {e}")

    # Export satellite POIs (category="satellite") separately
    try:
        all_pois = poi_manager.list_pois()
        satellite_pois = [
            poi for poi in all_pois if poi.category == "satellite"
        ]

        if satellite_pois:
            satellite_data = {
                "type": "satellite_pois",
                "pois": [poi.model_dump(mode="json") for poi in satellite_pois],
                "count": len(satellite_pois),
            }
            satellite_path = "pois/satellites.json"
            zf.writestr(
                satellite_path,
                json.dumps(satellite_data, indent=2, default=str),
            )
            manifest_files["pois"].append(satellite_path)
            logger.info(
                f"Added satellite POI data: {satellite_path} with {len(satellite_pois)} satellites"
            )
    except Exception as e:
        logger.error(f"Failed to add satellite POI data: {e}")


def _add_per_leg_exports_to_zip(
    zf: zipfile.ZipFile,
    mission: Mission,
    route_manager: RouteManager | None,
    poi_manager: POIManager | None,
    manifest_files: dict,
):
    """Generate and add per-leg exports (CSV, XLSX, PPTX, PDF) to zip archive.
    
    Args:
        zf: ZipFile to add files to
        mission: Mission object with legs
        route_manager: RouteManager instance
        poi_manager: POIManager instance
        manifest_files: Manifest dictionary to update
    """
    for leg in mission.legs:
        # Load timeline for this specific leg
        leg_timeline = load_mission_timeline(leg.id)
        if not leg_timeline:
            logger.warning(
                f"No timeline found for leg {leg.id}, skipping exports for this leg"
            )
            continue

        try:
            # CSV export
            csv_export = generate_timeline_export(
                export_format=TimelineExportFormat.CSV,
                mission=leg,  # Pass the leg as mission (MissionLeg has same fields)
                timeline=leg_timeline,
            )
            csv_path = f"exports/legs/{leg.id}/timeline.csv"
            zf.writestr(csv_path, csv_export.content)
            manifest_files["per_leg_exports"].append(csv_path)
            logger.info(f"Generated CSV for leg {leg.id}")
        except Exception as e:
            logger.error(f"Failed to generate CSV for leg {leg.id}: {e}")

        try:
            # XLSX export
            xlsx_export = generate_timeline_export(
                export_format=TimelineExportFormat.XLSX,
                mission=leg,
                timeline=leg_timeline,
                parent_mission_id=mission.id,
                route_manager=route_manager,
                poi_manager=poi_manager,
            )
            xlsx_path = f"exports/legs/{leg.id}/timeline.xlsx"
            zf.writestr(xlsx_path, xlsx_export.content)
            manifest_files["per_leg_exports"].append(xlsx_path)
            logger.info(f"Generated XLSX for leg {leg.id}")
        except Exception as e:
            logger.error(f"Failed to generate XLSX for leg {leg.id}: {e}")

        try:
            # PPTX export
            pptx_export = generate_timeline_export(
                export_format=TimelineExportFormat.PPTX,
                mission=leg,
                timeline=leg_timeline,
                parent_mission_id=mission.id,
                route_manager=route_manager,
                poi_manager=poi_manager,
            )
            pptx_path = f"exports/legs/{leg.id}/slides.pptx"
            zf.writestr(pptx_path, pptx_export.content)
            manifest_files["per_leg_exports"].append(pptx_path)
            logger.info(f"Generated PPTX for leg {leg.id}")
        except Exception as e:
            logger.error(f"Failed to generate PPTX for leg {leg.id}: {e}")

        try:
            # PDF export
            pdf_export = generate_timeline_export(
                export_format=TimelineExportFormat.PDF,
                mission=leg,
                timeline=leg_timeline,
                parent_mission_id=mission.id,
                route_manager=route_manager,
                poi_manager=poi_manager,
            )
            pdf_path = f"exports/legs/{leg.id}/report.pdf"
            zf.writestr(pdf_path, pdf_export.content)
            manifest_files["per_leg_exports"].append(pdf_path)
            logger.info(f"Generated PDF for leg {leg.id}")
        except Exception as e:
            logger.error(f"Failed to generate PDF for leg {leg.id}: {e}")


def _add_combined_mission_exports_to_zip(
    zf: zipfile.ZipFile,
    mission: Mission,
    route_manager: RouteManager | None,
    poi_manager: POIManager | None,
    manifest_files: dict,
):
    """Generate and add combined mission-level exports to zip archive.
    
    Args:
        zf: ZipFile to add files to
        mission: Mission object
        route_manager: RouteManager instance
        poi_manager: POIManager instance
        manifest_files: Manifest dictionary to update
    """
    try:
        logger.info(f"Generating combined mission-level exports for {mission.id}")

        # Combined CSV - stream to temp file
        with tempfile.NamedTemporaryFile(delete=True) as tmp_csv:
            generate_mission_combined_csv(mission, output_path=tmp_csv.name)
            zf.write(tmp_csv.name, "exports/mission/mission-timeline.csv")
            manifest_files["mission_exports"].append("exports/mission/mission-timeline.csv")

        # Combined XLSX - stream to temp file
        with tempfile.NamedTemporaryFile(delete=True, suffix=".xlsx") as tmp_xlsx:
            generate_mission_combined_xlsx(
                mission,
                route_manager=route_manager,
                poi_manager=poi_manager,
                output_path=tmp_xlsx.name,
            )
            zf.write(tmp_xlsx.name, "exports/mission/mission-timeline.xlsx")
            manifest_files["mission_exports"].append("exports/mission/mission-timeline.xlsx")

        with tempfile.NamedTemporaryFile(delete=True, suffix=".pptx") as tmp_pptx:
            generate_mission_combined_pptx(
                mission,
                route_manager=route_manager,
                poi_manager=poi_manager,
                output_path=tmp_pptx.name,
            )
            zf.write(tmp_pptx.name, "exports/mission/mission-slides.pptx")
            manifest_files["mission_exports"].append("exports/mission/mission-slides.pptx")

        with tempfile.NamedTemporaryFile(delete=True, suffix=".pdf") as tmp_pdf:
            generate_mission_combined_pdf(
                mission,
                route_manager=route_manager,
                poi_manager=poi_manager,
                output_path=tmp_pdf.name,
            )
            zf.write(tmp_pdf.name, "exports/mission/mission-report.pdf")
            manifest_files["mission_exports"].append("exports/mission/mission-report.pdf")

        logger.info("Combined mission-level exports complete")

    except Exception as e:
        logger.error(f"Failed to generate combined mission exports: {e}")


def _create_export_manifest(mission: Mission, manifest_files: dict) -> dict:
    """Create the manifest.json content for the export package.
    
    Args:
        mission: Mission object
        manifest_files: Dictionary with file lists by category
        
    Returns:
        Manifest dictionary
    """
    return {
        "version": "2.0",
        "mission_id": mission.id,
        "mission_name": mission.name,
        "mission_description": mission.description or "",
        "leg_count": len(mission.legs),
        "exported_at": datetime.now(timezone.utc).isoformat(),
        "file_structure": {
            "mission_data": manifest_files["mission_data"],
            "legs": manifest_files["legs"],
            "routes": manifest_files["routes"],
            "pois": manifest_files["pois"],
            "mission_exports": manifest_files["mission_exports"],
            "per_leg_exports": manifest_files["per_leg_exports"],
        },
        "statistics": {
            "total_files": (
                len(manifest_files["mission_data"])
                + len(manifest_files["legs"])
                + len(manifest_files["routes"])
                + len(manifest_files["pois"])
                + len(manifest_files["mission_exports"])
                + len(manifest_files["per_leg_exports"])
            ),
            "leg_files": len(manifest_files["legs"]),
            "route_files": len(manifest_files["routes"]),
            "poi_files": len(manifest_files["pois"]),
            "mission_export_files": len(manifest_files["mission_exports"]),
            "per_leg_export_files": len(manifest_files["per_leg_exports"]),
        },
    }


def export_mission_package(
    mission_id: str,
    route_manager: RouteManager,
    poi_manager: POIManager,
) -> IO[bytes]:
    """Export complete mission as zip archive.

    Package structure:
        mission-{id}.zip
        ├── mission.json
        ├── manifest.json
        ├── legs/
        │   ├── {leg-id-1}.json
        │   └── {leg-id-2}.json
        ├── routes/
        │   ├── {route-id-1}.kml
        │   └── {route-id-2}.kml
        ├── pois/
        │   ├── {leg-id-1}-pois.json
        │   └── {leg-id-2}-pois.json
        ├── exports/
        │   ├── mission/
        │   │   ├── mission-timeline.csv
        │   │   ├── mission-timeline.xlsx
        │   │   ├── mission-slides.pptx
        │   │   └── mission-report.pdf
        │   └── legs/
        │       ├── {leg-id-1}/
        │       │   ├── timeline.csv
        │       │   ├── timeline.xlsx
        │       │   ├── slides.pptx
        │       │   └── report.pdf
        │       └── {leg-id-2}/
        │           ├── timeline.csv
        │           ├── timeline.xlsx
        │           ├── slides.pptx
        │           └── report.pdf

    Args:
        mission_id: Mission to export
        route_manager: RouteManager instance for fetching route KML files
        poi_manager: POIManager instance for fetching mission POIs

    Returns:
        File-like object containing the zip archive. Caller must close it to delete the temp file.
    """
    mission = load_mission_v2(mission_id)

    if not mission:
        raise ExportPackageError(f"Mission {mission_id} not found")

    # Create a temporary file for the zip archive
    # delete=True ensures it's deleted when closed by the caller (StreamingResponse)
    zip_temp = tempfile.NamedTemporaryFile(delete=True)

    manifest_files = {
        "mission_data": ["mission.json", "manifest.json"],
        "legs": [],
        "routes": [],
        "pois": [],
        "mission_exports": [],
        "per_leg_exports": [],
    }

    try:
        with zipfile.ZipFile(zip_temp, "w", zipfile.ZIP_DEFLATED) as zf:
            # Add mission metadata and leg files
            _add_mission_metadata_to_zip(zf, mission, manifest_files)

            # Add route KML files
            _add_route_kmls_to_zip(zf, mission, route_manager, manifest_files)

            # Add POI data (leg-specific and satellites)
            _add_pois_to_zip(zf, mission, poi_manager, manifest_files)

            # Generate and add per-leg exports
            _add_per_leg_exports_to_zip(zf, mission, route_manager, poi_manager, manifest_files)

            # Generate and add combined mission-level exports
            _add_combined_mission_exports_to_zip(zf, mission, route_manager, poi_manager, manifest_files)

            # Create and add manifest
            manifest = _create_export_manifest(mission, manifest_files)
            manifest_json = json.dumps(manifest, indent=2)
            zf.writestr("manifest.json", manifest_json)
            logger.info(
                f"Created manifest.json with {manifest['statistics']['total_files']} files"
            )

    except Exception:
        # If anything fails, close and delete the temp file
        zip_temp.close()
        raise

    # Seek to beginning so it can be read
    zip_temp.seek(0)
    return zip_temp
