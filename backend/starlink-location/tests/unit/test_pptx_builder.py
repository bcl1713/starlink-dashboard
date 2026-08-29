"""Tests for PPTX slide generation with mission metadata."""

from unittest.mock import Mock, patch

from app.mission.exporter.__main__ import _cover_metadata_line
from app.mission.exporter.pptx_builder import _get_footer_metadata
from app.mission.exporter.pptx_styling import (
    STATUS_DEGRADED,
    TEXT_BLACK,
    add_status_badge,
)
from app.mission.models import Mission, MissionLeg, TimelineStatus, TransportConfig
from pptx import Presentation


def test_degraded_status_uses_yellow_with_black_text_palette():
    """Degraded PPT styling should be yellow, not orange, with black text."""
    assert (STATUS_DEGRADED[0], STATUS_DEGRADED[1], STATUS_DEGRADED[2]) == (
        255,
        255,
        0,
    )
    assert (TEXT_BLACK[0], TEXT_BLACK[1], TEXT_BLACK[2]) == (0, 0, 0)


def test_degraded_status_badge_uses_black_text_on_yellow_background():
    """Actual degraded status badges should render readable black text."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    badge = add_status_badge(slide, TimelineStatus.DEGRADED, 0, 0, 1, 0.3)
    paragraph = badge.text_frame.paragraphs[0]

    assert badge.fill.fore_color.rgb == STATUS_DEGRADED
    assert paragraph.font.color.rgb == TEXT_BLACK


def test_title_slide_cover_metadata_prefers_mission_revision_fields():
    """Cover metadata should not reuse stale Rev text from descriptions."""
    mission = Mission(
        id="mission-26-07",
        name="Mission 26-07",
        description="1 Leg | Rev 2",
        legs=[],
        metadata={"mission_number": "26-07", "revision": "5"},
    )

    assert _cover_metadata_line(mission, leg_count=1) == "1 Leg | Mission 26-07 | Rev 5"


class TestGetFooterMetadata:
    """Test footer metadata resolution for PPTX slides."""

    def test_with_none_mission(self):
        """Test that None mission returns empty string."""
        result = _get_footer_metadata(mission=None, parent_mission_id=None)
        assert result == ""

    def test_with_mission_leg_name_and_description(self):
        """Test footer with leg that has name and description."""
        leg = MissionLeg(
            id="leg-1",
            name="Leg 1",
            description="Test Flight",
            route_id="route-1",
            transports=TransportConfig(
                initial_x_satellite_id="X-1",
                initial_ka_satellite_ids=["Ka-1"],
            ),
        )
        result = _get_footer_metadata(mission=leg, parent_mission_id=None)
        assert result == "Leg 1 | Test Flight"

    def test_with_mission_leg_name_only(self):
        """Test footer with leg that has name but no description."""
        leg = MissionLeg(
            id="leg-1",
            name="Leg 1",
            description=None,
            route_id="route-1",
            transports=TransportConfig(
                initial_x_satellite_id="X-1",
                initial_ka_satellite_ids=["Ka-1"],
            ),
        )
        result = _get_footer_metadata(mission=leg, parent_mission_id=None)
        assert result == "Leg 1"

    def test_fallback_to_id_when_name_missing(self):
        """Test fallback to id when name is None."""
        # Create a mock object that doesn't have name attribute
        leg = Mock()
        leg.id = "leg-1"
        leg.name = None
        leg.description = "Test Flight"

        result = _get_footer_metadata(mission=leg, parent_mission_id=None)
        assert result == "leg-1 | Test Flight"

    @patch("app.mission.storage.load_mission_v2")
    def test_with_parent_mission_found(self, mock_load):
        """Test footer uses parent mission metadata when parent_mission_id provided."""
        # Setup parent mission
        parent = Mission(
            id="26-05",
            name="26-05",
            description="CONUS California",
            legs=[],
        )
        mock_load.return_value = parent

        # Setup leg
        leg = MissionLeg(
            id="leg-1",
            name="Leg 1 - Departure",
            description="First leg",
            route_id="route-1",
            transports=TransportConfig(
                initial_x_satellite_id="X-1",
                initial_ka_satellite_ids=["Ka-1"],
            ),
        )

        result = _get_footer_metadata(mission=leg, parent_mission_id="26-05")

        # Should use parent mission metadata
        assert result == "26-05 | CONUS California"
        mock_load.assert_called_once_with("26-05")

    @patch("app.mission.storage.load_mission_v2")
    def test_with_parent_mission_not_found(self, mock_load):
        """Test fallback to leg metadata when parent mission not found."""
        # Parent mission not found
        mock_load.return_value = None

        # Setup leg
        leg = MissionLeg(
            id="leg-1",
            name="Leg 1 - Departure",
            description="First leg",
            route_id="route-1",
            transports=TransportConfig(
                initial_x_satellite_id="X-1",
                initial_ka_satellite_ids=["Ka-1"],
            ),
        )

        result = _get_footer_metadata(mission=leg, parent_mission_id="26-05")

        # Should fall back to leg metadata
        assert result == "Leg 1 - Departure | First leg"
        mock_load.assert_called_once_with("26-05")

    @patch("app.mission.storage.load_mission_v2")
    def test_parent_mission_without_description(self, mock_load):
        """Test parent mission with name but no description."""
        # Setup parent mission without description
        parent = Mission(
            id="26-05",
            name="26-05",
            description=None,
            legs=[],
        )
        mock_load.return_value = parent

        # Setup leg
        leg = MissionLeg(
            id="leg-1",
            name="Leg 1",
            description="First leg",
            route_id="route-1",
            transports=TransportConfig(
                initial_x_satellite_id="X-1",
                initial_ka_satellite_ids=["Ka-1"],
            ),
        )

        result = _get_footer_metadata(mission=leg, parent_mission_id="26-05")

        # Should show parent name without separator
        assert result == "26-05"
        mock_load.assert_called_once_with("26-05")

    @patch("app.mission.storage.load_mission_v2")
    def test_parent_mission_fallback_to_id(self, mock_load):
        """Test parent mission falls back to id when name is None."""
        # Setup parent mission without name (using Mock since model doesn't allow None)
        parent = Mock()
        parent.id = "26-05"
        parent.name = None
        parent.description = "CONUS California"
        mock_load.return_value = parent

        # Setup leg
        leg = MissionLeg(
            id="leg-1",
            name="Leg 1",
            description="First leg",
            route_id="route-1",
            transports=TransportConfig(
                initial_x_satellite_id="X-1",
                initial_ka_satellite_ids=["Ka-1"],
            ),
        )

        result = _get_footer_metadata(mission=leg, parent_mission_id="26-05")

        # Should use parent id and description
        assert result == "26-05 | CONUS California"
        mock_load.assert_called_once_with("26-05")

    def test_mission_object_with_name_and_description(self):
        """Test with full Mission object (not just leg)."""
        mission = Mission(
            id="26-05",
            name="Mission 26-05",
            description="Full mission test",
            legs=[],
        )

        result = _get_footer_metadata(mission=mission, parent_mission_id=None)
        assert result == "Mission 26-05 | Full mission test"
